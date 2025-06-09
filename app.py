import streamlit as st
import os
import json
import time
import re
import hashlib
from io import BytesIO
import tempfile
from datetime import datetime, timezone
import traceback

from azure.storage.blob import BlobServiceClient
from azure.core.credentials import AzureKeyCredential
from azure.search.documents import SearchClient
from azure.search.documents.models import VectorizedQuery
from openai import AzureOpenAI
from dateutil import parser as date_parser
import pandas as pd

# Coqui TTS import (optional, ensure installed if used)
try:
    from TTS.api import TTS
    TTS_AVAILABLE = True
except ImportError:
    TTS_AVAILABLE = False
    # st.warning("TTS library not found. Text-to-speech for summaries will be unavailable.")
    # Warning will be shown in sidebar if summary button is disabled


st.set_page_config(layout="wide", page_title="AI Knowledge Pipeline")

# --- Environment Variable Configuration ---

# --- Azure OpenAI Service for RAG (Chat Model) ---
OPENAI_CHAT_API_KEY_RAG = os.getenv("OPENAI_CHAT_API_KEY_RAG", "YOUR_CHAT_API_KEY_FOR_RAG")
OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG = os.getenv("OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG", "YOUR_CHAT_DEPLOYMENT_NAME_FOR_RAG") # e.g., gpt-4, gpt-35-turbo-16k
OPENAI_CHAT_AZURE_ENDPOINT_RAG = os.getenv("OPENAI_CHAT_AZURE_ENDPOINT_RAG", "YOUR_CHAT_AZURE_OPENAI_ENDPOINT_FOR_RAG")
OPENAI_CHAT_API_VERSION_RAG = os.getenv("OPENAI_CHAT_API_VERSION_RAG", "2024-02-15-preview") # Or your preferred API version

# --- Azure OpenAI Service for Ingestion (Embedding Model) ---
OPENAI_EMBED_API_KEY_INGEST = os.getenv("OPENAI_EMBED_API_KEY_INGEST", "YOUR_EMBEDDING_API_KEY_FOR_INGESTION")
OPENAI_EMBED_LLM_DEPLOYMENT_NAME_INGEST = os.getenv("OPENAI_EMBED_LLM_DEPLOYMENT_NAME_INGEST", "YOUR_EMBEDDING_DEPLOYMENT_NAME_FOR_INGESTION") # e.g., text-embedding-ada-002
OPENAI_EMBED_AZURE_ENDPOINT_INGEST = os.getenv("OPENAI_EMBED_AZURE_ENDPOINT_INGEST", "YOUR_EMBEDDING_AZURE_OPENAI_ENDPOINT_FOR_INGESTION")
OPENAI_EMBED_API_VERSION_INGEST = os.getenv("OPENAI_EMBED_API_VERSION_INGEST", "2024-02-15-preview") # Or your preferred API version

# --- Azure Data Lake Storage (ADLS) Gen2 or Blob Storage ---
ADLS_ACCOUNT_URL = os.getenv("ADLS_ACCOUNT_URL", "YOUR_ADLS_BLOB_ACCOUNT_URL") # e.g., https://youraccount.blob.core.windows.net/
ADLS_CONTAINER_NAME = os.getenv("ADLS_CONTAINER_NAME", "YOUR_ADLS_CONTAINER_NAME")
STORAGE_ACCOUNT_KEY = os.getenv("STORAGE_ACCOUNT_KEY", "YOUR_ADLS_STORAGE_ACCOUNT_KEY")

# --- Azure AI Search ---
AZURE_SEARCH_ENDPOINT = os.getenv("AZURE_SEARCH_ENDPOINT", "YOUR_AZURE_AI_SEARCH_ENDPOINT") # e.g., https://yoursearchservice.search.windows.net
AZURE_SEARCH_API_KEY = os.getenv("AZURE_SEARCH_API_KEY", "YOUR_AZURE_AI_SEARCH_API_KEY")
AZURE_SEARCH_INDEX_NAME = os.getenv("AZURE_SEARCH_INDEX_NAME", "YOUR_AZURE_AI_SEARCH_INDEX_NAME")

# --- Text-to-Speech (TTS) Configuration ---
TTS_MODEL_NAME = os.getenv("TTS_MODEL_NAME", "tts_models/en/ljspeech/tacotron2-DDC")

# --- KT Specific Configurations ---
KT_UPLOAD_PATH_PREFIX = os.getenv("KT_UPLOAD_PATH_PREFIX", "user_kt_uploads") # Path prefix in ADLS for user-uploaded KT docs

# --- Application Target User/Context Configuration ---
# For scanning blobs related to a user (Tier 2 deep dive source)
APP_TARGET_USER_EMAIL_FOR_BLOBS = os.getenv("APP_TARGET_USER_EMAIL_FOR_BLOBS", "target_user_for_file_scan@example.com")
APP_TARGET_SOURCE_FOLDER_FOR_BLOBS = os.getenv("APP_TARGET_SOURCE_FOLDER_FOR_BLOBS", "Default_KT_Source_Folder_Path") # e.g., SharePoint_Data/user_specific_path_in_container

# For the specific KT session (interviewee, main subject of Q&A)
APP_KT_USER_EMAIL = os.getenv("APP_KT_USER_EMAIL", "kt_interviewee@example.com")
APP_KT_USER_NAME = os.getenv("APP_KT_USER_NAME", "KT Interviewee Name")
APP_KT_SOURCE_SYSTEM = os.getenv("APP_KT_SOURCE_SYSTEM", "Default_KT_Source_System_Name") # e.g., SharePoint_Data, SpecificProjectName


# --- Initialize Clients ---
try:
    chat_client = AzureOpenAI(
        api_key=OPENAI_CHAT_API_KEY_RAG,
        api_version=OPENAI_CHAT_API_VERSION_RAG,
        azure_endpoint=OPENAI_CHAT_AZURE_ENDPOINT_RAG,
    )
    embed_client = AzureOpenAI(
        api_key=OPENAI_EMBED_API_KEY_INGEST,
        api_version=OPENAI_EMBED_API_VERSION_INGEST,
        azure_endpoint=OPENAI_EMBED_AZURE_ENDPOINT_INGEST,
    )
    search_client = SearchClient(
        endpoint=AZURE_SEARCH_ENDPOINT,
        index_name=AZURE_SEARCH_INDEX_NAME,
        credential=AzureKeyCredential(AZURE_SEARCH_API_KEY),
    )
    blob_service_client = BlobServiceClient(account_url=ADLS_ACCOUNT_URL, credential=STORAGE_ACCOUNT_KEY)
except Exception as e:
    st.error(f"Fatal Error initializing Azure clients: {e}. Application cannot start. Check credentials or environment variables. Ensure placeholders are replaced with actual values.")
    st.stop()

# --- Dynamically construct filter string for Azure Search ---
# This filter is intended to retrieve content relevant to the KT session.
_filter_parts = []
_has_specific_kt_context = False

# Add filters if the values are not the default placeholders (implying they are intentionally set)
if APP_KT_USER_EMAIL and APP_KT_USER_EMAIL != "kt_interviewee@example.com":
    _filter_parts.append(f"employeeName eq '{APP_KT_USER_EMAIL}'")
    _has_specific_kt_context = True
if APP_KT_SOURCE_SYSTEM and APP_KT_SOURCE_SYSTEM != "Default_KT_Source_System_Name":
    _filter_parts.append(f"source eq '{APP_KT_SOURCE_SYSTEM}'")
    _has_specific_kt_context = True

# These are general KT-related source types, always include them in the OR condition.
_filter_parts.extend([
    "source eq 'UserUpload'",
    "source eq 'KTInterviewUpload'",
    "source eq 'InternalQASynthesis'"
])

KT_TARGET_RELEVANT_CONTENT_FILTER = " or ".join(_filter_parts)

if not _has_specific_kt_context:
    st.toast(
        "Context Warning: APP_KT_USER_EMAIL or APP_KT_SOURCE_SYSTEM not specifically set "
        "via environment variables. Search results may be less targeted.",
        icon="⚠️"
    )

# --- Utility Functions (largely unchanged, ensure they are complete from previous versions) ---
def non_null_list(value):
    if value is None: return []
    if isinstance(value, list): return value
    if isinstance(value, str): return [value] if value else []
    try: return list(value) if value else []
    except TypeError: return [value]

def sanitize_id_for_search(value: str) -> str:
    sanitized = re.sub(r'[^a-zA-Z0-9_=-]', '_', str(value))
    if len(sanitized) > 1024: sanitized = sanitized[:1024]
    if not sanitized:
        return "empty_id_fallback_" + hashlib.md5(str(value).encode('utf-8')).hexdigest()[:16]
    if sanitized.startswith("_"):
        sanitized = "id_" + sanitized
    return sanitized

def get_stable_id(blob_name: str, chunk_index: int) -> str:
    base_id_material = f"{blob_name}chunk{chunk_index}"
    hashed_suffix = hashlib.sha256(base_id_material.encode("utf-8")).hexdigest()[:16]
    sanitized_blob_prefix_parts = []
    for part in blob_name.split('/'):
        sanitized_part = re.sub(r'[^a-zA-Z0-9-]', '_', part)
        sanitized_part = re.sub(r'\+', '', sanitized_part).strip('_')
        if sanitized_part:
            sanitized_blob_prefix_parts.append(sanitized_part)
    prefix = "_".join(sanitized_blob_prefix_parts)
    max_prefix_len = 200
    prefix = prefix[:max_prefix_len]
    candidate_id = f"{prefix}{hashed_suffix}"
    return sanitize_id_for_search(candidate_id)

def remove_email_signature(text: str) -> str:
    if not text: return ""
    patterns = [
        r"(?i)(\bthanks?\b|\bregards?\b|\bcheers?\b|\bsincerely\b|\bkind regards\b|\bbest regards\b|\bbest\b)[\s,][\r\n]+.*$", # More greedy .*
        r"(?i)--\s*[\r\n]+.*$", r"(?i)^CONFIDENTIALITY NOTICE:.*$", r"(?i)^Disclaimer:.*$",
        r"(?i)^Sent from my .*", r"(?i)^Get Outlook for .*"
    ]
    cleaned_text = text
    for pattern in patterns:
        cleaned_text = re.sub(pattern, '', cleaned_text, flags=re.DOTALL | re.MULTILINE).strip()
    return cleaned_text


def semantic_chunk(text_content: str, max_chunk_len: int = 1500, overlap: int = 200) -> list[str]:
    if not isinstance(text_content, str) or not text_content.strip(): return []

    text_content = remove_email_signature(text_content)
    if not text_content.strip(): return []

    def get_sentences(paragraph):
        sentences = re.split(r'(?<=[.!?])\s+', paragraph.strip())
        return [s.strip() for s in sentences if s.strip()]

    chunks, current_chunk_parts, current_length = [], [], 0
    paragraphs = [p.strip() for p in text_content.split('\n\n') if p.strip()]

    if not paragraphs:
        if text_content.strip(): paragraphs = [text_content.strip()]
        else: return []

    for para_idx, paragraph in enumerate(paragraphs):
        if current_length + len(paragraph) + (len("\n\n") if current_chunk_parts else 0) <= max_chunk_len:
            current_chunk_parts.append(paragraph)
            current_length += len(paragraph) + (len("\n\n") if len(current_chunk_parts) > 1 else 0)
        else:
            if current_chunk_parts:
                chunks.append("\n\n".join(current_chunk_parts))
                current_chunk_parts, current_length = [], 0
                if overlap > 0 and chunks and chunks[-1]:
                    overlap_text_candidate = chunks[-1][-overlap:]
                    overlap_break = max(
                        overlap_text_candidate.rfind(". ") + 2,
                        overlap_text_candidate.rfind("\n\n") + 2,
                        overlap_text_candidate.rfind(" ") + 1,
                        0
                    )
                    actual_overlap_text = overlap_text_candidate[overlap_break:].strip()
                    if actual_overlap_text:
                        current_chunk_parts.append(actual_overlap_text)
                        current_length = len(actual_overlap_text)

            if len(paragraph) > max_chunk_len:
                sentences = get_sentences(paragraph)
                if not sentences:
                    prefix_for_long_para = ""
                    if current_chunk_parts:
                        prefix_for_long_para = "\n\n".join(current_chunk_parts) + "\n\n"
                        current_chunk_parts, current_length = [], 0

                    paragraph_to_split = prefix_for_long_para + paragraph
                    for i_fs in range(0, len(paragraph_to_split), max_chunk_len - overlap):
                        chunk_to_add_fs = paragraph_to_split[i_fs : i_fs + max_chunk_len]
                        chunks.append(chunk_to_add_fs.strip())
                    current_chunk_parts, current_length = [], 0
                    continue

                for sentence in sentences:
                    if current_length + len(sentence) + (len(" ") if current_chunk_parts else 0) <= max_chunk_len:
                        current_chunk_parts.append(sentence)
                        current_length += len(sentence) + (len(" ") if len(current_chunk_parts) > 1 else 0)
                    else:
                        if current_chunk_parts:
                            chunks.append(" ".join(current_chunk_parts))
                        current_chunk_parts, current_length = [], 0
                        if overlap > 0 and chunks and chunks[-1]:
                            overlap_text_candidate_sent = chunks[-1][-overlap:]
                            overlap_break_sent = max(
                                overlap_text_candidate_sent.rfind(". ") + 2,
                                overlap_text_candidate_sent.rfind(" ") + 1,
                                0
                            )
                            actual_overlap_text_sent = overlap_text_candidate_sent[overlap_break_sent:].strip()
                            if actual_overlap_text_sent:
                                current_chunk_parts.append(actual_overlap_text_sent)
                                current_length = len(actual_overlap_text_sent)
                        current_chunk_parts.append(sentence)
                        current_length += len(sentence)

                if current_chunk_parts:
                    chunks.append(" ".join(current_chunk_parts))
                    current_chunk_parts, current_length = [], 0
            else:
                current_chunk_parts = [paragraph]
                current_length = len(paragraph)

    if current_chunk_parts:
        chunks.append("\n\n".join(current_chunk_parts) if len(current_chunk_parts) > 1 and any("\n\n" in p for p in current_chunk_parts) else " ".join(current_chunk_parts))

    return [c.strip() for c in chunks if c.strip()]


def read_blob_to_memory(bsc: BlobServiceClient, container: str, blob: str) -> bytes:
    try:
        return bsc.get_blob_client(container, blob).download_blob().readall()
    except Exception as e:
        st.error(f"Read Blob Error for {blob}: {e}")
        raise

def extract_pptx_text(blob_data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        st.warning("python-pptx not installed. PPTX text extraction will fail.")
        return "Error: python-pptx not installed."
    try:
        prs = Presentation(BytesIO(blob_data))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text_runs.append("".join(run.text for run in para.runs))
            text_runs.append("\n<SLIDE_BREAK>\n")
        return "\n".join(text_runs).strip()
    except Exception as e:
        return f"Error extracting PPTX: {e}"

def extract_docx_text(blob_data: bytes) -> str:
    try:
        import docx
    except ImportError:
        st.warning("python-docx not installed. DOCX text extraction will fail.")
        return "Error: python-docx not installed."
    try:
        return "\n\n".join([p.text for p in docx.Document(BytesIO(blob_data)).paragraphs if p.text.strip()]).strip()
    except Exception as e:
        return f"Error extracting DOCX: {e}"

def extract_doc_text(blob_data: bytes, fn: str) -> str:
    try:
        import textract
    except ImportError:
        st.warning("textract not installed. DOC text extraction will fail.")
        return "Error: textract not installed."
    tmp_path = ""
    try:
        suffix = os.path.splitext(fn)[1] or ".doc"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            tmp.write(blob_data)
            tmp_path = tmp.name
        text = textract.process(tmp_path, encoding='utf-8').decode('utf-8', errors='ignore').strip()
        return text
    except Exception as e:
        return f"Error extracting .doc with textract for '{fn}': {e}"
    finally:
        if tmp_path and os.path.exists(tmp_path):
            try: os.remove(tmp_path)
            except Exception: pass

def extract_pdf_text_and_images(blob_data: bytes, blob_name: str) -> list[dict]:
    extracted = []
    try:
        import fitz  # PyMuPDF
    except ImportError:
        st.warning("PyMuPDF (fitz) not installed. PDF text extraction will fail.")
        extracted.append({"type": "text", "content": "Error: PyMuPDF (fitz) not installed."})
        return extracted
    try:
        doc = fitz.open(stream=blob_data, filetype="pdf")
        full_text = ""
        for page_num, page in enumerate(doc):
            full_text += page.get_text("text", sort=True) + "\n\n"
        if full_text.strip():
            extracted.append({"type": "text", "content": full_text.strip()})
        else:
            extracted.append({"type": "text", "content": f"No text could be extracted from PDF {blob_name}."})
        doc.close()
    except Exception as e:
        extracted.append({"type": "text", "content": f"Error processing PDF {blob_name} with PyMuPDF: {e}"})
    return extracted

def extract_excel_text(blob_data: bytes, ext: str) -> str:
    engine = 'openpyxl' if ext == '.xlsx' else 'xlrd' if ext == '.xls' else None
    if not engine: return f"Error: Unsupported Excel ext: {ext}"
    try:
        xls = pd.ExcelFile(BytesIO(blob_data), engine=engine)
        parts = []
        for sheet_name in xls.sheet_names:
            parts.append(f"\n--- Sheet: {sheet_name} ---\n")
            df = xls.parse(sheet_name, header=None)
            if not df.empty:
                parts.append('\n'.join([' '.join(map(str, row_tuple)) for row_tuple in df.fillna('').itertuples(index=False)]))
            else: parts.append("(Sheet is empty)")
            parts.append("\n")
        txt = "".join(parts).strip()
        return txt if txt else "(Excel file contains no text data)"
    except ImportError:
        msg = f"{engine} library not installed for {ext}."
        st.error(msg); return f"Error: {msg}"
    except Exception as e:
        st.error(f"Error processing Excel: {e}"); return f"Error processing Excel: {e}"

def extract_json_from_llm_response(response_text: str) -> dict | list | None:
    if not response_text or not isinstance(response_text, str):
        return {"error": "Empty/invalid LLM response"}

    match_markdown = re.search(r"```json\s*([\s\S]*?)\s*```", response_text, re.DOTALL)
    if match_markdown: json_str = match_markdown.group(1).strip()
    else:
        start_brace = response_text.find('{'); start_bracket = response_text.find('[')
        if start_brace == -1 and start_bracket == -1: json_str = response_text
        elif start_brace != -1 and (start_bracket == -1 or start_brace < start_bracket):
            end_brace = response_text.rfind('}')
            if end_brace > start_brace: json_str = response_text[start_brace : end_brace+1]
            else: json_str = response_text
        elif start_bracket != -1 and (start_brace == -1 or start_bracket < start_brace):
            end_bracket = response_text.rfind(']')
            if end_bracket > start_bracket: json_str = response_text[start_bracket : end_bracket+1]
            else: json_str = response_text
        else: json_str = response_text

    try: return json.loads(json_str)
    except json.JSONDecodeError as e_initial:
        obj_match = re.search(r"(\{[\s\S]*\})", response_text)
        list_match = re.search(r"(\[[\s\S]*\])", response_text)
        json_like_str_candidate = None
        if obj_match and list_match: json_like_str_candidate = obj_match.group(0) if obj_match.start() < list_match.start() else list_match.group(0)
        elif obj_match: json_like_str_candidate = obj_match.group(0)
        elif list_match: json_like_str_candidate = list_match.group(0)

        if json_like_str_candidate:
            try: return json.loads(json_like_str_candidate)
            except json.JSONDecodeError as e_sub:
                error_detail = f"Initial parse failed ({e_initial}). Substring JSON parse error: {e_sub}."
                original_sample = response_text[:200] + "..." + response_text[-200:] if len(response_text) > 400 else response_text
                return {"error": f"{error_detail} Original text sample: {original_sample}"}
        return {"error": f"No JSON object/list found after multiple attempts. Initial error: {e_initial}. Original text sample: {response_text[:200]}"}


def format_date_for_azure_search(date_input) -> str | None:
    if not date_input: return None
    dt = None
    if isinstance(date_input, datetime): dt = date_input
    elif isinstance(date_input, str):
        try: dt = date_parser.parse(date_input)
        except (ValueError, TypeError, OverflowError):
            if re.match(r"^\d{4}-\d{2}-\d{2}T\d{2}:\d{2}:\d{2}(\.\d+)?$", date_input): return date_input + "Z"
            if re.match(r"^\d{4}-\d{2}-\d{2}$", date_input): return date_input + "T00:00:00Z"
            st.warning(f"Could not parse date string: {date_input}"); return None
    else: st.warning(f"Unsupported date input type: {type(date_input)} for value {date_input}"); return None

    if dt.tzinfo is None or dt.tzinfo.utcoffset(dt) is None: dt = dt.replace(tzinfo=timezone.utc)
    else: dt = dt.astimezone(timezone.utc)
    return dt.isoformat().replace("+00:00", "Z")

def transcribe_audio_with_whisper(audio_blob_data: bytes, original_filename: str) -> str:
    st.info(f"Audio transcription for {original_filename} is a placeholder. Phase 3 will implement this.")
    return "Placeholder: Audio transcription not yet implemented."
# --- END Utility Functions ---

# --- Global Constants and Field Definitions ---
ALL_AZURE_SEARCH_FIELDS = [
    {"name": "id", "type": "Edm.String", "key": True, "filterable": True, "sortable": True, "facetable": False, "searchable": False},
    {"name": "chunkText", "type": "Edm.String", "searchable": True, "retrievable": True},
    {"name": "vector", "type": "Collection(Edm.Single)", "dimensions": 1536, "vectorSearchProfile": "default-vector-profile"},
    {"name": "employeeId", "type": "Edm.String", "searchable": True, "filterable": True, "facetable": True},
    {"name": "employeeName", "type": "Edm.String", "searchable": True, "filterable": True, "facetable": True},
    {"name": "roleTitle", "type": "Edm.String", "searchable": True, "filterable": True, "facetable": True},
    {"name": "department", "type": "Edm.String", "searchable": True, "filterable": True, "facetable": True},
    {"name": "businessUnit", "type": "Edm.String", "searchable": True, "filterable": True, "facetable": True},
    {"name": "tenure", "type": "Edm.String", "searchable": True, "filterable": True, "facetable": True},
    {"name": "lastWorkingDay", "type": "Edm.DateTimeOffset", "filterable": True, "sortable": True, "facetable": True},
    {"name": "primaryLocation", "type": "Edm.String", "searchable": True, "filterable": True, "facetable": True},
    {"name": "documentType", "type": "Edm.String", "searchable": True, "filterable": True, "facetable": True},
    {"name": "sourceFileName", "type": "Edm.String", "filterable": True, "facetable": True, "sortable": True, "searchable": True},
    {"name": "sourceFilePath", "type": "Edm.String", "filterable": True, "retrievable": True},
    {"name": "extension", "type": "Edm.String", "filterable": True, "facetable": True},
    {"name": "timestamp", "type": "Edm.DateTimeOffset", "filterable": True, "sortable": True, "facetable": True},
    {"name": "chunkIndex", "type": "Edm.Int32", "filterable": True, "sortable": True},
    {"name": "projectName", "type": "Edm.String", "searchable": True, "filterable": True, "facetable": True},
    {"name": "projectDescription", "type": "Edm.String", "searchable": True},
    {"name": "projectRole", "type": "Edm.String", "searchable": True, "filterable": True},
    {"name": "projectDuration", "type": "Edm.String", "searchable": True, "filterable": True},
    {"name": "projectStatus", "type": "Edm.String", "searchable": True, "filterable": True},
    {"name": "projectStartDate", "type": "Edm.DateTimeOffset", "filterable": True, "sortable": True},
    {"name": "projectEndDate", "type": "Edm.DateTimeOffset", "filterable": True, "sortable": True},
    {"name": "projectDeliverables", "type": "Collection(Edm.String)", "searchable": True, "filterable": True, "facetable": True},
    {"name": "projectMilestones", "type": "Collection(Edm.ComplexType)", "fields": [
        {"name": "date", "type": "Edm.String"},
        {"name": "milestone", "type": "Edm.String", "searchable": True}
    ]},
    {"name": "projectDataSources", "type": "Collection(Edm.ComplexType)", "fields": [
        {"name": "system", "type": "Edm.String", "searchable": True, "filterable": True},
        {"name": "tables_files", "type": "Collection(Edm.String)", "searchable": True, "filterable": True}
    ]},
    {"name": "projectProcessSteps", "type": "Collection(Edm.ComplexType)", "fields": [
        {"name": "step", "type": "Edm.String", "searchable": True},
        {"name": "tool", "type": "Edm.String", "searchable": True, "filterable": True}
    ]},
    {"name": "projectStakeholders", "type": "Collection(Edm.ComplexType)", "fields": [
        {"name": "name", "type": "Edm.String", "searchable": True, "filterable": True},
        {"name": "role", "type": "Edm.String", "searchable": True},
        {"name": "responsibility", "type": "Edm.String", "searchable": True}
    ]},
    {"name": "projectChallenges", "type": "Collection(Edm.ComplexType)", "fields": [
        {"name": "issue", "type": "Edm.String", "searchable": True},
        {"name": "resolution", "type": "Edm.String", "searchable": True}
    ]},
     {"name": "projectKPIs", "type": "Collection(Edm.ComplexType)", "fields": [
        {"name": "metric", "type": "Edm.String", "searchable": True, "filterable": True},
        {"name": "target_value", "type": "Edm.String", "searchable": True},
        {"name": "business_action", "type": "Edm.String", "searchable": True}
    ]},
    {"name": "projectOpenItems", "type": "Collection(Edm.ComplexType)", "fields": [
        {"name": "item", "type": "Edm.String", "searchable": True},
        {"name": "risk", "type": "Edm.String", "filterable": True},
        {"name": "next_step", "type": "Edm.String", "searchable": True}
    ]},
    {"name": "achievementTitle", "type": "Edm.String", "searchable": True, "filterable": True, "facetable": True},
    {"name": "achievementDescription", "type": "Edm.String", "searchable": True},
    {"name": "achievementDate", "type": "Edm.DateTimeOffset", "filterable": True, "sortable": True},
    {"name": "domainName", "type": "Edm.String", "searchable": True, "filterable": True, "facetable": True},
    {"name": "insightDescription", "type": "Edm.String", "searchable": True},
    {"name": "kpiName", "type": "Collection(Edm.String)", "searchable": True, "filterable": True, "facetable": True},
    {"name": "kpiValue", "type": "Edm.Double", "filterable": True, "facetable": True, "sortable": True},
    {"name": "kpiPeriod", "type": "Edm.String", "filterable": True, "facetable": True, "sortable": True},
    {"name": "skillName", "type": "Collection(Edm.String)", "searchable": True, "filterable": True, "facetable": True},
    {"name": "skillLevel", "type": "Collection(Edm.String)", "searchable": True, "filterable": True, "facetable": True},
    {"name": "toolName", "type": "Collection(Edm.String)", "searchable": True, "filterable": True, "facetable": True},
    {"name": "toolCategory", "type": "Collection(Edm.String)", "searchable": True, "filterable": True, "facetable": True},
    {"name": "contactName", "type": "Collection(Edm.String)", "searchable": True, "filterable": True, "facetable": True},
    {"name": "contactRole", "type": "Collection(Edm.String)", "searchable": True, "filterable": True, "facetable": True},
    {"name": "contactPurpose", "type": "Collection(Edm.String)", "searchable": True},
    {"name": "customMetadata", "type": "Collection(Edm.ComplexType)", "fields": [
        {"name": "key", "type": "Edm.String", "searchable": True, "filterable": True},
        {"name": "value", "type": "Edm.String", "searchable": True, "filterable": True}
    ]},
    {"name": "uploaded_by", "type": "Edm.String", "filterable": True, "facetable": True, "searchable": True},
    {"name": "upload_date", "type": "Edm.DateTimeOffset", "filterable": True, "sortable": True, "facetable": True},
    {"name": "source", "type": "Edm.String", "filterable": True, "facetable": True, "searchable": True},
    {"name": "tags", "type": "Collection(Edm.String)", "filterable": True, "facetable": True, "searchable": True},
    {"name": "description", "type": "Edm.String", "searchable": True},
    {"name": "version", "type": "Edm.String", "filterable": True, "facetable": True},
    {"name": "kt_session_id", "type": "Edm.String", "filterable": True, "facetable": True},
    {"name": "kt_session_question", "type": "Edm.String", "searchable": True, "filterable": True, "retrievable": True},
    {"name": "kt_session_answer", "type": "Edm.String", "searchable": True, "filterable": True, "retrievable": True}
]
KNOWN_EXTENSIONLESS_TEXT_FILES = {"metadata", "readme", "license", "config", "description", "gitignore", "dockerfile", "makefile"}
DEFAULT_RAG_RETRIEVAL_FIELDS = list(set([
    "id", "chunkText", "sourceFileName", "documentType", "employeeName",
    "projectName", "projectRole", "projectDeliverables", "skillName", "timestamp", "uploaded_by",
    "source", "tags", "department", "domainName", "roleTitle",
    "kt_session_question", "kt_session_answer"
]))
# --- END Global Constants ---

# --- Structured KT Interview Question Generation (CogniLink Persona) ---
FEW_SHOT_PROMPT_BLOCK_V2_ADAPTED = """
You are CogniLink, an advanced, empathetic, and dynamic AI interviewer. Your primary goal is to conduct a comprehensive and natural-flowing Knowledge Transfer (KT) interview. You are patient, encouraging, and aim to understand the user's contributions, responsibilities, and knowledge deeply. You MUST adapt your questioning based on the user's responses AND the system's validation feedback provided in the conversation history.

INTERVIEW FLOW:
The interview generally follows this structure:
1. Introduction: Welcome and initial project identification.
2. Deep Dive per Project: For each identified project, thoroughly explore various facets:
    - Objectives & User's Role/Contributions
    - Key Deliverables
    - Major Milestones & Timeline
    - Data Sources & Systems Used
    - Key Process Steps & Tools
    - Stakeholders (internal/external) & Collaboration
    - Challenges Faced & Solutions Implemented
    - Key Performance Indicators (KPIs) & Metrics
    - Open Items, Risks, or Pending Tasks
3. Project Transition: Smoothly move to the next project once one is comprehensively covered.
4. Wrap-up: Concluding questions about overall learnings, advice for successors, or any missed critical information.

ADVANCED CHAIN OF THOUGHT (CoT) - Follow these 4 steps BEFORE generating any question:
1. REVIEW HISTORY & VALIDATION: What was the last AI question? What was the user's answer? Critically, what was the system's validation score and feedback for that answer (e.g., `System Validation of User's Last Answer: Score=0.3 (Source: User Upload: report.docx). Details: The document mentions timelines but not the specific budget figures asked for.`)?
2. ASSESS CURRENT STATE: Which project are we currently deep-diving into? Which KT topic (from the list above) for THIS project should be next? Are there gaps indicated by a low validation score (e.g., below 0.6) or specific feedback in the system validation message that require a clarifying sub-question BEFORE moving to a new topic?
3. DETERMINE NEXT LOGICAL STEP:
    - If validation was low or indicates missing details for the current topic, ask a targeted sub-question to clarify or get more specific information related to the *last question asked*.
    - If the current topic is well-covered (good validation), move to the next logical KT topic for the *current project*.
    - If all topics for the current project are covered, transition to the next identified project or ask if there are more projects.
    - If all projects are covered, move to wrap-up questions.
    - If the user explicitly states "next", conclude the current deep dive if appropriate and move to the next main topic or project based on the overall interview flow.
4. FORMULATE QUESTION: Craft the single, most appropriate next question as a JSON object with "question" and "response_structure" fields. Ensure the `response_structure` provides clear examples for the user.

FEW-SHOT EXAMPLES:

EXAMPLE 0 (Introduction - First Question):
System: (Initial call, no conversation history)
CogniLink (CoT):
1. Review: No history.
2. Assess: Start of interview. Need to identify projects.
3. Next Step: Ask for an overview of primary projects.
4. Formulate:
JSON:
{
  "question": "Welcome! I'm CogniLink, here to help with your knowledge transfer. To begin, could you list the primary projects you've been significantly involved with recently? For each, please briefly mention your main role.",
  "response_structure": {
    "projects_overview": [
      {"project_name": "e.g., Project Phoenix", "your_role": "e.g., Lead Analyst"},
      {"project_name": "e.g., Client Integration Gamma", "your_role": "e.g., Project Manager"}
    ]
  }
}

EXAMPLE 1 (Project Identification & Basic Details):
System: (User lists Project Alpha - Lead, Project Beta - Analyst)
User: "I worked on Project Alpha as Lead Developer and Project Beta as a supporting Data Analyst."
CogniLink (CoT):
1. Review: User listed Alpha (Lead), Beta (Analyst). No validation yet.
2. Assess: Projects identified. Need more details for Alpha first. Duration and status are good next points.
3. Next Step: Ask for duration and status for Project Alpha.
4. Formulate:
JSON:
{
  "question": "Thanks! Let's start with 'Project Alpha'. Could you provide the approximate start and end dates (or duration) and its current completion status (e.g., Completed, In-progress, On Hold)?",
  "response_structure": {
    "project_name_context": "Project Alpha",
    "start_date": "YYYY-MM-DD",
    "end_date": "YYYY-MM-DD (or 'Ongoing')",
    "status": "e.g., Completed"
  }
}

EXAMPLE 2 (Deep Dive - Project Objectives & Role - First topic for Project Alpha):
System: (User provided duration/status for Project Alpha: Jan 2023 – Dec 2023, Completed)
User: "Project Alpha ran from Jan 2023 to Dec 2023 and it's completed."
System Validation of User's Last Answer: Score=1.0 (Source: User Self-Attested). Details: User provided clear dates and status.
CogniLink (CoT):
1. Review: Project Alpha details: Jan 2023-Dec 2023, Completed. Validation is high.
2. Assess: Current project is Alpha. Next topic: objective and user's specific contributions.
3. Next Step: Ask about Project Alpha's main objective and user's role specifics.
4. Formulate:
JSON:
{
  "question": "Great. Focusing on 'Project Alpha', what was its main objective or the problem it aimed to solve? And, as the Lead Developer, what were your specific responsibilities and key contributions?",
  "response_structure": {
    "project_name_context": "Project Alpha",
    "objective": "e.g., To refactor the legacy billing system to improve performance by 30% and reduce errors.",
    "my_role_and_contributions": "e.g., I led the architectural redesign, developed the core processing modules using Python and Spark, and coordinated the work of 3 junior developers."
  }
}

EXAMPLE 2.1 (Sub-question/Clarification due to Low Validation):
System: (User answered about Project Alpha objectives, but vaguely)
User: "Project Alpha was about making the system better."
System Validation of User's Last Answer: Score=0.3 (Source: RAG from 'project_charter.pdf'). Details: RAG context mentions performance and scalability, but user's answer is too general and doesn't reflect these specifics.
CogniLink (CoT):
1. Review: User answer "making system better" for Alpha objectives. Validation score is low (0.3). System feedback suggests charter mentioned performance/scalability.
2. Assess: Still on Project Alpha, topic: objectives. Low validation indicates need for clarification.
3. Next Step: Ask a sub-question to get more specific details about the objectives, perhaps hinting at performance/scalability if appropriate.
4. Formulate:
JSON:
{
  "question": "Could you elaborate a bit more on 'making the system better' for Project Alpha? For instance, were there specific aspects like performance improvements, new functionalities, or cost reductions that were targeted?",
  "response_structure": {
    "project_name_context": "Project Alpha",
    "specific_objectives_clarification": "e.g., Yes, the main goal was to improve data processing speed by at least 50% and integrate new regulatory reporting features."
  }
}

EXAMPLE 3 (Deep Dive - Deliverables for Project Alpha):
System: (User provided clear objectives and contributions for Project Alpha)
User: "The objective of Alpha was to cut processing time by 50%. I designed the new data pipeline and led the coding team."
System Validation of User's Last Answer: Score=0.9 (Source: User Upload 'Alpha_Role_Desc.docx'). Details: Document confirms user's leadership and focus on data pipeline.
CogniLink (CoT):
1. Review: Objectives/contributions for Alpha are clear. Validation is high.
2. Assess: Current project: Alpha. Next topic: deliverables.
3. Next Step: Ask for main deliverables of Project Alpha.
4. Formulate:
JSON:
{
  "question": "That's very clear, thank you. What were the main deliverables or tangible outputs produced for 'Project Alpha'?",
  "response_structure": {
    "project_name_context": "Project Alpha",
    "deliverables": [
      "e.g., The new data processing pipeline software",
      "e.g., Technical architecture and design documents",
      "e.g., Deployment scripts and operational runbooks"
    ]
  }
}

EXAMPLE 4 (Deep Dive - Milestones for Project Alpha):
System: (User provided deliverables for Project Alpha)
User: (Lists deliverables for Alpha)
System Validation of User's Last Answer: Score=0.8 (Source: User Self-Attested). Details: User listed relevant deliverables.
CogniLink (CoT):
1. Review: Deliverables for Alpha obtained. Good validation.
2. Assess: Current project: Alpha. Next topic: milestones and timeline.
3. Next Step: Ask for major milestones and overall timeline for Project Alpha.
4. Formulate:
JSON:
{
  "question": "Excellent. Could you outline the major milestones and the overall timeline (confirming start/end dates if they differ from overall project duration) for 'Project Alpha'?",
  "response_structure": {
    "project_name_context": "Project Alpha",
    "start_date_confirmed": "YYYY-MM-DD",
    "end_date_confirmed": "YYYY-MM-DD",
    "milestones": [
      {"date": "YYYY-MM-DD", "milestone_description": "e.g., Requirements finalized and signed off"},
      {"date": "YYYY-MM-DD", "milestone_description": "e.g., Core engine development complete (MVP)"},
      {"date": "YYYY-MM-DD", "milestone_description": "e.g., UAT completed and approved"},
      {"date": "YYYY-MM-DD", "milestone_description": "e.g., Production deployment and go-live"}
    ]
  }
}

EXAMPLE 5 (Deep Dive - Data Sources for Project Alpha):
System: (User provided milestones for Project Alpha)
User: (Lists milestones for Alpha)
System Validation of User's Last Answer: Score=0.9 (Source: User Upload 'Alpha_Timeline.pptx'). Details: Presentation slides confirm milestones.
CogniLink (CoT):
1. Review: Milestones for Alpha obtained. High validation.
2. Assess: Current project: Alpha. Next topic: data sources/systems.
3. Next Step: Ask about key data sources for Project Alpha.
4. Formulate:
JSON:
{
  "question": "Very helpful. For 'Project Alpha', what were the key systems, applications, or data sources (e.g., specific databases like Oracle/SQL Server, APIs, Kafka topics, shared drives with Excel/CSV files) that you regularly interacted with or relied upon? Please specify table names or crucial file names if they were central.",
  "response_structure": {
    "project_name_context": "Project Alpha",
    "data_sources": [
      {"system_or_source_type": "e.g., Oracle Database", "details": "e.g., Customer_Master table, Orders table"},
      {"system_or_source_type": "e.g., Internal REST API", "details": "e.g., ProductDetailsService API endpoint"},
      {"system_or_source_type": "e.g., S3 Bucket", "details": "e.g., Raw_Sensor_Data folder, daily CSV uploads"}
    ]
  }
}

EXAMPLE 6 (Deep Dive - Process Steps & Tools for Project Alpha):
System: (User provided data sources for Project Alpha)
User: (Lists data sources for Alpha)
CogniLink (CoT):
1. Review: Data sources for Alpha obtained. Assume good validation.
2. Assess: Current project: Alpha. Next topic: key process steps and tools used.
3. Next Step: Ask about the main process steps and tools for Project Alpha.
4. Formulate:
JSON:
{
  "question": "Understanding the data flow is key. Now, could you walk me through the main process steps involved in 'Project Alpha' from start to finish, and mention any specific software, tools, or programming languages that were critical for these steps?",
  "response_structure": {
    "project_name_context": "Project Alpha",
    "process_steps_tools": [
      {"step_description": "e.g., Data ingestion from S3", "tools_languages_used": "e.g., Python (Boto3), Airflow"},
      {"step_description": "e.g., Data transformation and cleaning", "tools_languages_used": "e.g., Apache Spark (PySpark), Pandas"},
      {"step_description": "e.g., Loading data into Oracle DWH", "tools_languages_used": "e.g., SQL, Python (cx_Oracle)"},
      {"step_description": "e.g., API development for data exposure", "tools_languages_used": "e.g., FastAPI (Python), Docker"}
    ]
  }
}

EXAMPLE 7 (Deep Dive - Stakeholders for Project Alpha):
System: (User provided process steps for Project Alpha)
User: (Lists process steps for Alpha)
CogniLink (CoT):
1. Review: Process steps for Alpha obtained. Assume good validation.
2. Assess: Current project: Alpha. Next topic: stakeholders.
3. Next Step: Ask about key stakeholders for Project Alpha.
4. Formulate:
JSON:
{
  "question": "That gives a good overview of the workflow. Who were the key stakeholders (e.g., business units, specific individuals, external clients) for 'Project Alpha'? What was their involvement or what did you primarily interact with them about?",
  "response_structure": {
    "project_name_context": "Project Alpha",
    "stakeholders": [
      {"name_or_group": "e.g., Sales Department", "role_or_interest": "e.g., Primary users of the new system reports", "interaction_nature": "e.g., Requirements gathering, UAT feedback"},
      {"name_or_group": "e.g., Jane Doe (Head of Operations)", "role_or_interest": "e.g., Project Sponsor", "interaction_nature": "e.g., Weekly status updates, decision approvals"},
      {"name_or_group": "e.g., External Vendor X", "role_or_interest": "e.g., Provided a key software component", "interaction_nature": "e.g., Integration support, issue resolution"}
    ]
  }
}

EXAMPLE 8 (Deep Dive - Challenges for Project Alpha):
System: (User provided stakeholders for Project Alpha)
User: (Lists stakeholders for Alpha)
CogniLink (CoT):
1. Review: Stakeholders for Alpha obtained. Assume good validation.
2. Assess: Current project: Alpha. Next topic: challenges.
3. Next Step: Ask about challenges and solutions for Project Alpha.
4. Formulate:
JSON:
{
  "question": "Understanding the people involved is crucial. What were some of the significant challenges or obstacles you encountered during 'Project Alpha', and how did you or the team address them?",
  "response_structure": {
    "project_name_context": "Project Alpha",
    "challenges_solutions": [
      {"challenge": "e.g., Unexpectedly poor quality of legacy data", "solution_or_mitigation": "e.g., Implemented an extensive data cleansing pre-processing stage and worked with data owners to correct sources."},
      {"challenge": "e.g., Scope creep from a key stakeholder", "solution_or_mitigation": "e.g., Established a formal change request process and negotiated phasing for non-critical new features."}
    ]
  }
}

EXAMPLE 9 (Deep Dive - KPIs for Project Alpha):
System: (User provided challenges for Project Alpha)
User: (Lists challenges for Alpha)
CogniLink (CoT):
1. Review: Challenges for Alpha obtained. Assume good validation.
2. Assess: Current project: Alpha. Next topic: KPIs/metrics.
3. Next Step: Ask about how success was measured for Project Alpha.
4. Formulate:
JSON:
{
  "question": "It's good to know how those were handled. How was the success of 'Project Alpha' measured? Were there specific Key Performance Indicators (KPIs) or metrics that were tracked?",
  "response_structure": {
    "project_name_context": "Project Alpha",
    "kpis_metrics": [
      {"kpi_or_metric": "e.g., Data processing time reduction", "target_value": "e.g., Target 50% reduction", "actual_outcome": "e.g., Achieved 55% reduction"},
      {"kpi_or_metric": "e.g., System error rate", "target_value": "e.g., Target < 0.1% errors", "actual_outcome": "e.g., Achieved 0.05% error rate"},
      {"kpi_or_metric": "e.g., User adoption rate", "target_value": "e.g., 90% of Sales team using new reports within 3 months", "actual_outcome": "e.g., 95% adoption achieved"}
    ]
  }
}

EXAMPLE 9.1 (Deep Dive - Open Items for Project Alpha - Last topic for Alpha):
System: (User provided KPIs for Project Alpha)
User: (Lists KPIs for Alpha)
CogniLink (CoT):
1. Review: KPIs for Alpha obtained. Assume good validation.
2. Assess: Current project: Alpha. Next topic: open items/pending tasks. This is likely the last deep dive topic for Alpha.
3. Next Step: Ask about any open items for Project Alpha.
4. Formulate:
JSON:
{
  "question": "Those are solid results for 'Project Alpha'. Were there any open items, known issues, planned future enhancements, or pending tasks related to Project Alpha when you concluded your involvement or as of now?",
  "response_structure": {
    "project_name_context": "Project Alpha",
    "open_items": [
      {"item_description": "e.g., Phase 2 performance tuning for Q3 peak load", "current_status_or_owner": "e.g., Scheduled for Q3, Ops team to monitor", "associated_risk_if_any": "e.g., Minor risk of slowdown if not addressed before peak"},
      {"item_description": "e.g., Documentation for new API endpoints needs final review", "current_status_or_owner": "e.g., Draft complete, assigned to Sarah for review", "associated_risk_if_any": "e.g., None critical"}
    ]
  }
}


EXAMPLE 10 (Project Transition - All topics for Project Alpha covered, moving to Project Beta):
System: (User provided open items for Project Alpha. All deep dive topics for Alpha are now covered.)
User: (Lists open items for Alpha)
System Validation of User's Last Answer: Score=0.9 (Source: User Self-Attested). Details: User provided relevant open items.
CogniLink (CoT):
1. Review: Open items for Alpha covered. Validation good. All topics for Project Alpha seem covered. User previously mentioned Project Beta.
2. Assess: Project Alpha deep dive complete.
3. Next Step: Transition to Project Beta. Start with its duration and status, similar to how Alpha was started.
4. Formulate:
JSON:
{
  "question": "Thank you for the thorough details on Project Alpha! Now, let's shift to 'Project Beta', where you mentioned you were a supporting Data Analyst. Could you provide its approximate start/end dates and current status?",
  "response_structure": {
    "project_name_context": "Project Beta",
    "start_date": "YYYY-MM-DD",
    "end_date": "YYYY-MM-DD (or 'Ongoing')",
    "status": "e.g., In-progress"
  }
}
# (The interview would then proceed to deep dive into Project Beta using similar questions as for Project Alpha: Objectives, Role, Deliverables, etc.)

EXAMPLE 15 (Wrap-up - All identified projects covered):
System: (All topics for all identified projects like Alpha, Beta, etc., have been discussed)
User: (Provides final details for the last project)
System Validation of User's Last Answer: Score=1.0. Details: User confirms completion.
CogniLink (CoT):
1. Review: All projects covered. Last response validated well.
2. Assess: End of project-specific deep dives. Time for wrap-up.
3. Next Step: Ask general concluding questions.
4. Formulate:
JSON:
{
  "question": "We've covered a lot of ground on your projects, which is incredibly helpful! Just a couple of final questions: Are there any critical pieces of knowledge, key contacts, or general advice you'd like to pass on to someone taking over your responsibilities? Or perhaps any general skills or tools you found particularly vital in your role that we haven't touched upon?",
  "response_structure": {
    "overall_key_knowledge_or_contacts": "e.g., Always double-check the data reconciliation scripts after any upstream changes. John Smith from IT is the go-to for any DWH access issues.",
    "general_advice_for_successor": "e.g., Get familiar with the Python data science stack (Pandas, NumPy, Scikit-learn) as it's heavily used.",
    "other_critical_info_or_skills": "e.g., Strong SQL skills are essential. Also, an understanding of Agile methodologies was beneficial."
  }
}

EXAMPLE 16 (End of Interview - After Wrap-up):
System: (User has answered wrap-up questions.)
User: "I think that covers most of it. Good luck to my successor!"
System Validation of User's Last Answer: Score=1.0 (Source: User Self-Attested). Details: User indicates they are done.
CogniLink (CoT):
1. Review: User has answered wrap-up and indicated completion.
2. Assess: Interview is complete.
3. Next Step: Formally conclude the interview.
4. Formulate:
JSON:
{
  "question": "This has been an extremely productive session. Thank you so much for your time and detailed responses! This information will be invaluable. That concludes our KT interview. [END_OF_INTERVIEW]",
  "response_structure": {
    "acknowledgement": "e.g., Thank you / Sounds good."
  }
}


INSTRUCTION:
You MUST use your Advanced Chain of Thought (CoT) before every question:
1. **Review History & Validation:** Carefully examine the conversation history. Pay close attention to the user's last answer and, most importantly, any `System Validation of User's Last Answer` messages. These messages contain a `Score` (0.0 to 1.0) and `Details` that indicate how well the user's answer was substantiated.
2. **Assess Current State:** Determine which project is currently under discussion. Identify the next logical KT topic for this project (Objectives, Deliverables, etc.). If a System Validation score for the previous answer on the current topic was low (e.g., < 0.6) or the details indicate missing information, a sub-question to clarify that specific point is necessary before moving on.
3. **Determine Next Logical Step:** Based on the Interview Flow, your CoT, and the validation feedback:
    - Ask a **sub-question** if clarification is needed for the current topic due to low validation.
    - Move to the **next KT topic** for the current project if the current one is well-covered.
    - **Transition to a new project** if all topics for the current project are complete.
    - Move to **wrap-up questions** if all projects are covered.
    - If the user says "next", interpret this as a signal to advance, potentially concluding the current deep dive and moving to the next major topic/project.
    - If the interview is logically complete (all projects and wrap-up done), generate the "END_OF_INTERVIEW" question.
4. **Formulate Question:** Provide your response ONLY as a JSON object with "question" and "response_structure" keys. The `question` should be natural and empathetic. The `response_structure` should guide the user by providing examples of the information expected.

Your goal is to conduct a comprehensive, empathetic, and dynamic KT interview. Adapt your questions based on the flow of conversation and the validation insights.
Focus on one project at a time, cycling through all relevant KT topics for that project before moving to the next.
Ensure all questions are polite and encouraging.

Now, based on the user profile below and the conversation history provided (which may include system validation messages), generate the *next single appropriate KT question*.
"""

def get_user_profile_text_for_kt(employee_name, employee_email, known_info_summary=None):
    profile = f"User Profile (Person providing KT):\nName: {employee_name}\nEmail: {employee_email}\n"
    if known_info_summary:
        profile += f"Existing Information Summary (from RAG about this user, use as context if relevant, but prioritize direct interview answers):\n{known_info_summary}\n"
    else:
        profile += "This is the beginning of the KT session. Little prior specific project detail about this user is confirmed unless stated in conversation history.\n"
    return profile

def generate_single_kt_interview_question(
    employee_name: str,
    employee_email: str,
    conversation_history: list,
    known_info_summary: str = None
):
    global chat_client, OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG

    user_profile_text = get_user_profile_text_for_kt(employee_name, employee_email, known_info_summary)

    # Use the new adapted prompt
    messages = [{"role": "system", "content": FEW_SHOT_PROMPT_BLOCK_V2_ADAPTED + user_profile_text}]

    for turn in conversation_history:
        content = turn["content"]
        # Ensure content is string, especially for system observations which might be dicts
        if not isinstance(content, str):
            try:
                # If it's a dict (like our validation message) or list, dump to JSON string
                if isinstance(content, (dict, list)):
                    content_str = json.dumps(content)
                else: # For other non-string types, just convert to string
                    content_str = str(content)
            except TypeError: # Fallback if json.dumps fails for some reason
                content_str = str(content)
        else:
            content_str = content
        messages.append({"role": turn["role"], "content": content_str})

    st.session_state.debug_llm_messages_kt_qgen = messages

    try:
        response = chat_client.chat.completions.create(
            model=OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG,
            messages=messages,
            temperature=0.2, # Slightly lower for more focused questions based on structured CoT
            max_tokens=1500, # Increased for potentially complex response_structures and CoT in examples
        )
        content_str = response.choices[0].message.content

        parsed_json = extract_json_from_llm_response(content_str)

        if isinstance(parsed_json, dict) and "question" in parsed_json and "response_structure" in parsed_json:
            return parsed_json
        else:
            st.error(f"CogniLink QGen: LLM did not return valid JSON. Raw: {content_str[:500]}. Parsed: {parsed_json}")
            st.session_state.debug_kt_qgen_error_response = {"raw": content_str, "parsed": parsed_json}
            # Fallback question if LLM fails to produce valid JSON
            return {
                "question": "I had a slight hiccup formulating the next structured question. Can you tell me what you'd like to cover next, or perhaps reiterate your last point? You can also type 'next' to move to a new topic.",
                "response_structure": {"generic_response": "Your text answer here, or type 'next'"}
            }
    except Exception as e:
        st.error(f"CogniLink QGen LLM Error: {e}")
        traceback.print_exc()
        return {
            "question": "An error occurred while generating the next question. Please try again, provide a general update on your projects, or type 'next' to move on.",
            "response_structure": {"error_response": "Details about the error or your next point, or type 'next'"}
        }
# --- END Structured KT Interview Question Generation ---


# --- RAG, Validation, Ingestion, Tier 2 (largely unchanged, ensure they are complete) ---
@st.cache_data(ttl=600)
def retrieve_top_contexts_for_rag(query: str, filter_expression: str = None,
                                  top_k: int = 3, select_fields: list = None) -> list[dict]:
    global embed_client, search_client, OPENAI_EMBED_LLM_DEPLOYMENT_NAME_INGEST, DEFAULT_RAG_RETRIEVAL_FIELDS
    if not query:
        st.warning("RAG: Empty query provided.")
        return []

    fields_to_select = select_fields[:] if select_fields is not None else DEFAULT_RAG_RETRIEVAL_FIELDS[:]
    if "chunkText" not in fields_to_select and "kt_session_answer" not in fields_to_select :
        fields_to_select.insert(0, "chunkText")
    if "id" not in fields_to_select:
        fields_to_select.append("id")
    fields_to_select = list(dict.fromkeys(fields_to_select))

    try:
        query_vector_response = embed_client.embeddings.create(
            input=[query],
            model=OPENAI_EMBED_LLM_DEPLOYMENT_NAME_INGEST
        )
        if not query_vector_response.data or not query_vector_response.data[0].embedding:
            st.error("RAG Embedding Error: Received no embedding vector from API.")
            traceback.print_exc()
            return []
        query_vector = query_vector_response.data[0].embedding
    except Exception as e:
        st.error(f"RAG Embedding Error: {e}")
        traceback.print_exc()
        return []

    vector_query = VectorizedQuery(vector=query_vector, k_nearest_neighbors=top_k, fields="vector", exhaustive=True)
    try:
        results = search_client.search(
            search_text=query,
            vector_queries=[vector_query],
            filter=filter_expression,
            select=fields_to_select,
            top=top_k
        )
        docs = []
        for res in results:
            doc = {field: res.get(field) for field in fields_to_select if field in res}
            doc["@search.score"] = res.get("@search.score", 0.0)
            docs.append(doc)

        if any(d.get("@search.score", 0.0) > 0 for d in docs):
             docs.sort(key=lambda x: x.get("@search.score", 0.0), reverse=True)

        if not docs:
            st.warning(f"RAG Search: No documents found for query '{query[:50]}...' with filter '{filter_expression}'. The index might be empty, the filter too restrictive, or the query too dissimilar to content.")
        return docs
    except Exception as e:
        st.error(f"RAG Search Error for query '{query[:50]}...': {e}")
        traceback.print_exc()
        return []

@st.cache_data(ttl=3600)
def generate_ai_questions_from_indexed_data(
    focus_entity_name: str, entity_type_description: str,
    num_questions: int, use_rag_summary_context: bool,
    search_filter_for_summary: str = None
) -> list[str]:
    global chat_client, OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG, DEFAULT_RAG_RETRIEVAL_FIELDS

    initial_context_summary_str = "No initial summary context was retrieved or requested for question generation."
    if use_rag_summary_context:
        try:
            summary_query_for_rag = (
                f"Summarize key information, projects, achievements, and responsibilities for "
                f"'{focus_entity_name}' based on their documents to help formulate targeted questions about this entity's work."
            )
            summary_contexts = retrieve_top_contexts_for_rag(
                query=summary_query_for_rag,
                filter_expression=search_filter_for_summary,
                top_k=5,
                select_fields=DEFAULT_RAG_RETRIEVAL_FIELDS
            )
            if summary_contexts:
                summary_parts = []
                for i, ctx in enumerate(summary_contexts):
                    text_content = ctx.get('chunkText') or ctx.get('kt_session_answer')
                    source_info = ctx.get('sourceFileName', 'Unknown Source')
                    if text_content:
                        summary_parts.append(
                           f"Context Snippet {i+1} (from '{source_info}'):\n"
                           f"{text_content[:350]}...\n"
                        )
                if summary_parts:
                    initial_context_summary_str = "\n---\n".join(summary_parts)
                else:
                     initial_context_summary_str = "Retrieved context chunks did not yield a textual summary from chunkText or kt_session_answer."
            else:
                initial_context_summary_str = (
                    f"No relevant context documents found via RAG for '{focus_entity_name}' "
                    f"(filter: {search_filter_for_summary}) for question generation."
                )
        except Exception as e_rag:
            initial_context_summary_str = f"Error retrieving RAG summary context for QGen: {e_rag}"
            st.warning(initial_context_summary_str)
            traceback.print_exc()

    system_prompt_qgen = (
        "You are an expert AI assistant specializing in knowledge extraction and inquiry. "
        f"Your task is to generate insightful and specific questions about '{focus_entity_name}' "
        "based on potential indexed knowledge from their documents. "
        "The questions should be suitable for a knowledge transfer session, aiming to uncover detailed and actionable information."
    )
    user_prompt_parts_qgen = [
        f"Generate exactly {num_questions} distinct, insightful questions regarding: "
        f"'{entity_type_description}' concerning '{focus_entity_name}'.",
        "The questions should probe for details about this entity's work, including (but not limited to): "
        "Specific projects involved in (names, objectives, timelines, status, role, contributions), "
        "Key achievements and outcomes, Skills utilized or developed, "
        "Significant KPIs tracked or influenced (names, values, periods, trends), "
        "Important contacts, collaborators, or vendors worked with, "
        "Domains or business units supported, Challenges faced and solutions implemented, "
        "Processes, tools, or data sources frequently used or developed, "
        "Key learnings, insights, or assumptions from their work.",
        "Frame questions to be open-ended to encourage detailed responses."
    ]
    if use_rag_summary_context and "Error" not in initial_context_summary_str \
       and "No relevant" not in initial_context_summary_str \
       and "Retrieved context chunks did not yield" not in initial_context_summary_str:
        user_prompt_parts_qgen.append(
            f"\nTo help you formulate highly relevant and targeted questions, consider this summary "
            f"derived from documents related to '{focus_entity_name}':\n--- Summary of Available Information ---\n"
            f"{initial_context_summary_str}\n--- End Summary ---"
        )
    else:
        user_prompt_parts_qgen.append(
            f"\nNo specific pre-summary context is available for '{focus_entity_name}'. "
            "Generate general probing questions based on the typical knowledge areas for an entity like this, "
            "focusing on the elements listed above."
        )
    user_prompt_parts_qgen.append(
        f"\nFormat your response as a JSON object with a single key 'questions', where the value is a list of the generated question strings. For example: {{\"questions\": [\"What were the primary objectives of Project X?\", \"Can you describe your role in the Q1 marketing campaign?\"]}}"
    )
    final_user_prompt_qgen = "\n".join(user_prompt_parts_qgen)

    try:
        response = chat_client.chat.completions.create(
            model=OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG,
            messages=[
                {"role": "system", "content": system_prompt_qgen},
                {"role": "user", "content": final_user_prompt_qgen}
            ],
            temperature=0.6,
            max_tokens=200 * num_questions + 400,
            response_format={"type": "json_object"}
        )
        content = response.choices[0].message.content
        parsed_data = extract_json_from_llm_response(content)
        questions = []
        if isinstance(parsed_data, dict) and "error" not in parsed_data:
            if "questions" in parsed_data and isinstance(parsed_data["questions"], list):
                questions = [str(q).strip() for q in parsed_data["questions"] if isinstance(q, str) and q.strip()]
            else:
                for key, value in parsed_data.items():
                    if isinstance(value, list) and all(isinstance(item, str) for item in value):
                        questions = [str(q).strip() for q in value if q.strip()]
                        st.info(f"QGen (Automated): Found questions under key '{key}' instead of 'questions'.")
                        break
        elif isinstance(parsed_data, list) and all(isinstance(item, str) for item in parsed_data):
            questions = [str(q).strip() for q in parsed_data if q.strip()]
            st.info("QGen (Automated): LLM returned a direct list of questions.")

        if not questions:
            st.warning(f"QGen (Automated) LLM unexpected format or empty list. Raw: {content[:300]}. Parsed: {parsed_data}")
            return []
        return list(dict.fromkeys(questions))[:num_questions]
    except Exception as e:
        st.error(f"QGen (Automated) LLM Error: {e}")
        traceback.print_exc()
        return []

def answer_question_with_rag(
    question_to_answer: str,
    search_filter_expression: str = None,
    top_k_contexts: int = 3,
    context_select_fields: list = None,
    max_context_length_for_llm: int = 7000,
    conversation_history: list[dict] = None
) -> tuple[str, bool, list[dict]]:
    global chat_client, OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG, DEFAULT_RAG_RETRIEVAL_FIELDS
    if not question_to_answer:
        return "No question provided.", False, []

    rag_select_fields = context_select_fields if context_select_fields else DEFAULT_RAG_RETRIEVAL_FIELDS
    if "chunkText" not in rag_select_fields and "kt_session_answer" not in rag_select_fields:
        rag_select_fields.append("chunkText")
    rag_select_fields = list(dict.fromkeys(rag_select_fields + ["id", "sourceFileName", "documentType", "projectName", "kt_session_question"]))


    retrieved_contexts = retrieve_top_contexts_for_rag(
        question_to_answer, search_filter_expression, top_k_contexts, rag_select_fields
    )
    if not retrieved_contexts:
        return "No relevant RAG contexts found in the Azure Search Index for this specific question.", False, []

    parts, current_total_length = [], 0
    used_context_ids = set()
    for i, doc in enumerate(retrieved_contexts):
        if doc.get('id') in used_context_ids:
            continue

        context_text_to_use = doc.get("chunkText", "")
        source_prefix = "Document Snippet"
        if not context_text_to_use and doc.get("kt_session_answer"):
            context_text_to_use = f"Previously answered question: '{doc.get('kt_session_question', '')}' Answer: '{doc.get('kt_session_answer', '')}'"
            source_prefix = "Stored Q&A"

        if not context_text_to_use:
            continue

        src_fn = doc.get("sourceFileName", "Unknown Source")
        doc_type = doc.get("documentType", "")
        proj_name = doc.get("projectName", "")
        source_details = f"Source: {src_fn}"
        if doc_type: source_details += f" (Type: {doc_type})"
        if proj_name: source_details += f" (Project: {proj_name})"

        header = f"{source_prefix} {i+1} (ID: {doc.get('id','N/A')}, {source_details}):\n{context_text_to_use}\n"

        if current_total_length + len(header) > max_context_length_for_llm and parts:
            st.info(f"RAG (Answer): Context truncated at {max_context_length_for_llm} chars for LLM. Used {len(parts)} of {len(retrieved_contexts)} retrieved docs.")
            break
        parts.append(header)
        current_total_length += len(header)
        used_context_ids.add(doc.get('id'))

    if not parts:
        return "Retrieved RAG contexts, but could not prepare any suitable text snippets for the LLM.", True, retrieved_contexts

    context_string_for_llm = "\n---\n".join(parts)

    messages_for_llm = []
    system_prompt_text = (
        "You are a highly analytical AI assistant. Your task is to answer the user's 'Current Question' with high precision, "
        "strictly and solely based on the information available in the 'Provided Context Chunks'. "
        "Consider the 'Previous Conversation History' for context, especially if the current question is a follow-up. "
        "Do not use any external knowledge. If the information is not in the contexts, clearly state that. Be factual and concise."
    )
    messages_for_llm.append({"role": "system", "content": system_prompt_text})

    if conversation_history:
        knowledge_history = [
            msg for msg in conversation_history
            if msg.get("tier") != "Conversational" and msg.get("role") in ["user", "assistant"]
        ]
        relevant_history = knowledge_history[-4:]
        if relevant_history:
            messages_for_llm.append({"role": "user", "content": "--- Previous Conversation History (for context) ---"})
            for hist_msg in relevant_history:
                messages_for_llm.append({"role": hist_msg["role"], "content": hist_msg["content"]})
            messages_for_llm.append({"role": "user", "content": "--- End of Previous Conversation History ---"})

    user_prompt_text = (
        f"*Provided Context Chunks (for the Current Question):*\n{context_string_for_llm}\n\n\n"
        f"*Current Question from User:*\n{question_to_answer}\n\n"
        "Based *only* on the 'Provided Context Chunks' and being mindful of the 'Previous Conversation History' (if any), what is the answer to the 'Current Question from User'?\nAnswer:"
    )
    messages_for_llm.append({"role": "user", "content": user_prompt_text})

    try:
        response = chat_client.chat.completions.create(
            model=OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG,
            messages=messages_for_llm,
            temperature=0.0,
            max_tokens=1200
        )
        llm_answer = response.choices[0].message.content.strip()
        return llm_answer, True, retrieved_contexts
    except Exception as e:
        st.error(f"Tier 1 Answer LLM Error (with history): {e}")
        traceback.print_exc()
        return "Error synthesizing Tier 1 answer due to an LLM issue.", True, retrieved_contexts

DESIGNATED_NO_ANSWER_PHRASE = "NO_SPECIFIC_ANSWER_FOUND_IN_THIS_SNIPPET"
TIER2_CHUNK_FEW_SHOTS = f"""
    Example 1:
    Original Question: "What were the start and end dates for Project Phoenix?"
    Text Snippet from document 'project_plan.docx': "Project Phoenix was initiated on Jan 15, 2023. Phase 1 is planned to conclude by June 30, 2023, with subsequent phases lasting an additional 6 months."
    Answer: "Project Phoenix was initiated on Jan 15, 2023. Phase 1 concludes June 30, 2023. Subsequent phases add 6 months. Overall explicit end date for entire project not in this snippet."

    Example 2:
    Original Question: "What was the allocated budget for the Project Fusion marketing campaign?"
    Text Snippet from document 'budget_fy23.xlsx': "Fusion - R&D: $150,000; Fusion - Travel: $20,000. Trinity integration: $50,000. Mars initiative marketing: $75,000."
    Answer: "{DESIGNATED_NO_ANSWER_PHRASE}"

    Example 3:
    Original Question: "Who was the main point of contact for Vendor X?"
    Text Snippet from document 'meeting_notes_vendor_x.txt': "Discussion with Vendor X team (John Doe, Jane Smith). Key decisions made by our lead, Mark Johnson."
    Answer: "The snippet mentions John Doe and Jane Smith from Vendor X, and Mark Johnson as our lead for decisions. It does not explicitly state who the 'main point of contact' for Vendor X was."
    """
UNSATISFACTORY_PHRASES = [
    "information not found", "could not find specific information",
    "context does not specify", "context does not provide",
    "not available in the provided documents", "i could not find specific information",
    "not available in the provided context", "information is not available in the provided documents",
    "no specific information to answer this question", "unable to answer",
    "cannot answer based on the provided context",
    DESIGNATED_NO_ANSWER_PHRASE.lower(),
    "no answer in snippet", "does not mention any specific", "could not find information",
    "does not contain sufficient information", "does not detail", "i'm sorry", "i cannot",
    "i do not have enough information", "the provided text does not contain",
    "the document does not say", "the context doesn't mention", "no specific answer was found"
]
TIER2_CHUNK_SIZE = 2000
TIER2_CHUNK_OVERLAP = 300
MAX_CHUNKS_PER_DOC_TIER2 = 10

def is_answer_unsatisfactory(answer: str) -> bool:
    if not answer or not answer.strip():
        return True
    answer_lower = answer.lower()
    for phrase in UNSATISFACTORY_PHRASES:
        if phrase in answer_lower:
            return True
    if len(answer.split()) < 4 and ("yes" not in answer_lower and "no" not in answer_lower):
        if not any(char.isdigit() for char in answer):
            return True
    return False

def perform_tier2_deep_dive_for_one_question(
    unanswered_question: str, files_to_search: list,
    chat_client_instance, blob_service_client_instance, adls_container: str
) -> list[dict]:
    st.write(f"  Tier 2 Deep Dive: Analyzing up to {len(files_to_search)} files for question: \"{unanswered_question[:70]}...\"")
    all_doc_answers_status = []

    for blob_item_data in files_to_search:
        blob_name = blob_item_data["name"]
        blob_size = blob_item_data.get("size", 0)
        base_filename = os.path.basename(blob_name)

        if blob_size == 0:
            st.caption(f"    Tier 2: Skipping empty file {base_filename}.")
            all_doc_answers_status.append({
                "document_name": blob_name, "answer": None, "status": "skipped_empty_file"
            })
            continue

        st.caption(f"    Tier 2: Processing {base_filename} ({blob_size / 1024:.1f} KB)...")

        full_text_content_for_doc = ""
        try:
            blob_data_bytes = read_blob_to_memory(blob_service_client_instance, adls_container, blob_name)
            file_extension = os.path.splitext(blob_name)[-1].lower()

            extractor_map = {
                '.pptx': extract_pptx_text,
                '.docx': extract_docx_text,
            }

            if file_extension in extractor_map:
                full_text_content_for_doc = extractor_map[file_extension](blob_data_bytes)
            elif file_extension == '.doc':
                 full_text_content_for_doc = extract_doc_text(blob_data_bytes, blob_name)
            elif file_extension == '.pdf':
                pdf_parts = extract_pdf_text_and_images(blob_data_bytes, blob_name)
                full_text_content_for_doc = "\n\n".join([
                    p['content'] for p in pdf_parts
                    if p['type'] == 'text' and p.get('content') and "Error:" not in p['content'] and "No text could be extracted" not in p['content']
                ])
            elif file_extension in ['.xlsx', '.xls']:
                 full_text_content_for_doc = extract_excel_text(blob_data_bytes, file_extension)
            elif file_extension in [
                '.txt', '.json', '.html', '.htm', '.csv', '.md', '.xml', '.yaml', '.yml', '.log', '.sql',
                '.py', '.js', '.java', '.c', '.cpp', '.h', '.cs', '.rb', '.php', '.go', '.rs', '.swift',
                '.sh', '.ps1', '.ini', '.cfg', '.conf', '.rtf', '.sample', ''
            ] or (not file_extension and base_filename.lower() in KNOWN_EXTENSIONLESS_TEXT_FILES):
                decoded_text = blob_data_bytes.decode('utf-8', errors='ignore')
                if file_extension in ['.html', '.htm']:
                    text_for_llm = re.sub(r'<style(?:\s[^>]*)?>.*?</style>', '', decoded_text, flags=re.DOTALL | re.IGNORECASE)
                    text_for_llm = re.sub(r'<script(?:\s[^>]*)?>.*?</script>', '', text_for_llm, flags=re.DOTALL | re.IGNORECASE)
                    text_for_llm = re.sub(r'<[^>]+>', ' ', text_for_llm)
                    full_text_content_for_doc = re.sub(r'\s+', ' ', text_for_llm).strip()
                else:
                    full_text_content_for_doc = decoded_text
                full_text_content_for_doc = remove_email_signature(full_text_content_for_doc)
            else:
                all_doc_answers_status.append({
                    "document_name": blob_name, "answer": None, "status": f"unsupported_file_type_tier2 ({file_extension})"
                })
                continue

            if "Error:" in full_text_content_for_doc or not full_text_content_for_doc.strip():
                status_msg = "parser_error_or_empty_content"
                if "Error: python-" in full_text_content_for_doc or "Error: textract" in full_text_content_for_doc or "Error: PyMuPDF" in full_text_content_for_doc:
                    status_msg = "parser_dependency_missing"
                all_doc_answers_status.append({
                    "document_name": blob_name, "answer": None, "status": status_msg
                })
                continue
        except Exception as e_parse:
            st.error(f"Tier 2 Parsing/Reading Error for {blob_name}: {e_parse}")
            traceback.print_exc()
            all_doc_answers_status.append({
                "document_name": blob_name, "answer": None, "status": "critical_parsing_error"
            })
            continue

        tier2_semantic_chunks = semantic_chunk(full_text_content_for_doc, TIER2_CHUNK_SIZE, TIER2_CHUNK_OVERLAP)
        if not tier2_semantic_chunks:
            all_doc_answers_status.append({
                "document_name": blob_name, "answer": None, "status": "no_tier2_chunks_from_content"
            })
            continue

        positive_answer_snippets_from_doc = []
        system_prompt_chunk_qa = (
            "You are a meticulous AI information retriever. Your task is to determine if the 'Text Snippet' (extracted from the document "
            f"'{base_filename}') contains a direct and clear answer to the 'Original Question'.\n"
            "Follow these steps carefully:\n"
            "1. Deeply understand the 'Original Question'. What specific information is it seeking?\n"
            "2. Thoroughly read the 'Text Snippet'.\n"
            "3. Compare the information in the snippet against the question.\n"
            "4. If the snippet directly and unambiguously answers the question, extract the relevant sentence(s) or phrase(s) VERBATIM or as a very close paraphrase. Your answer should be self-contained.\n"
            "5. If the snippet provides a partial answer, a strong clue, or contextually relevant information that helps address the question (even if not fully), extract that information.\n"
            "6. If the snippet contains NO information relevant to answering the question, you MUST respond ONLY with the exact phrase: " + DESIGNATED_NO_ANSWER_PHRASE + "\n"
            "Do NOT invent information or infer beyond what is explicitly stated in the snippet.\n"
            "Here are some examples of how to respond:\n" + TIER2_CHUNK_FEW_SHOTS +
            "\nNow, evaluate the following:\n"
        )

        processed_chunk_count_for_doc = 0
        for chunk_idx, chunk_text in enumerate(tier2_semantic_chunks):
            if processed_chunk_count_for_doc >= MAX_CHUNKS_PER_DOC_TIER2:
                st.caption(f"    Tier 2: {base_filename} - Reached max {MAX_CHUNKS_PER_DOC_TIER2} chunks for this doc.")
                break
            if not chunk_text.strip():
                continue

            user_prompt_chunk_qa = (
                f"Original Question: \"{unanswered_question}\"\n\n"
                f"Text Snippet from document '{base_filename}' (Chunk {chunk_idx+1} of {len(tier2_semantic_chunks)}):\n"
                f"```text\n{chunk_text}\n```\n\n"
                "Based *only* on the content of this specific 'Text Snippet', does it provide an answer (full or partial) to the 'Original Question'? "
                "Respond with the extracted answer or the designated no-answer phrase.\nAnswer:"
            )
            try:
                response_chunk_qa = chat_client_instance.chat.completions.create(
                    model=OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG,
                    messages=[
                        {"role": "system", "content": system_prompt_chunk_qa},
                        {"role": "user", "content": user_prompt_chunk_qa}
                    ],
                    temperature=0.0,
                    max_tokens=800
                )
                extracted_answer_from_chunk = response_chunk_qa.choices[0].message.content.strip()

                if extracted_answer_from_chunk and extracted_answer_from_chunk != DESIGNATED_NO_ANSWER_PHRASE:
                    if not is_answer_unsatisfactory(extracted_answer_from_chunk):
                        positive_answer_snippets_from_doc.append({
                            "answer": extracted_answer_from_chunk,
                            "source_doc": blob_name,
                            "chunk_index": chunk_idx
                        })
            except Exception as e_llm_chunk:
                st.error(f"Tier 2 LLM Error (Doc: {base_filename}, Chunk {chunk_idx+1}): {e_llm_chunk}")

            processed_chunk_count_for_doc +=1

        if positive_answer_snippets_from_doc:
            all_doc_answers_status.extend(positive_answer_snippets_from_doc)
        else:
             all_doc_answers_status.append({
                "document_name": blob_name,
                "answer": None,
                "status": "no_answer_snippets_found_in_doc_chunks_tier2"
            })

    final_positive_snippets_for_synthesis = [
        item for item in all_doc_answers_status if item.get("answer") is not None
    ]
    return final_positive_snippets_for_synthesis

def synthesize_answers_from_snippets(
    question: str,
    answer_snippets_with_sources: list[dict],
    chat_client_instance,
    conversation_history: list[dict] = None
) -> str:
    if not answer_snippets_with_sources:
        return f"No specific answer snippets were found that directly address the question: \"{question}\""

    if len(answer_snippets_with_sources) == 1:
        return f"{answer_snippets_with_sources[0]['answer']} (Source: {os.path.basename(answer_snippets_with_sources[0]['source_doc'])})"

    snippet_text_for_llm = ""
    for i, item in enumerate(answer_snippets_with_sources):
        snippet_text_for_llm += (f"Snippet {i+1} from document '{os.path.basename(item['source_doc'])}':\n"
                                 f"{item['answer']}\n---\n")

    messages_for_llm = []
    system_prompt_text = (
        "You are an expert AI assistant tasked with synthesizing information from multiple text snippets into a single, coherent answer to the 'Original Question'. "
        "Consider the 'Previous Conversation History' for context if provided. Base your answer *only* on the provided snippets and history. "
        "If snippets conflict, note the conflict. If the information is still insufficient after reviewing all snippets, clearly state what's missing. Preserve important details from the snippets."
    )
    messages_for_llm.append({"role": "system", "content": system_prompt_text})

    if conversation_history:
        knowledge_history = [msg for msg in conversation_history if msg.get("tier") != "Conversational" and msg.get("role") in ["user", "assistant"]]
        relevant_history = knowledge_history[-4:]
        if relevant_history:
            messages_for_llm.append({"role": "user", "content": "--- Previous Conversation History (for context) ---"})
            for hist_msg in relevant_history:
                messages_for_llm.append({"role": hist_msg["role"], "content": hist_msg["content"]})
            messages_for_llm.append({"role": "user", "content": "--- End of Previous Conversation History ---"})

    user_prompt_text = (
        f"Original question: \"{question}\"\n\n"
        f"Snippets for synthesis:\n{snippet_text_for_llm}\n\n"
        "Please synthesize these snippets into a single, comprehensive answer for the 'Original question', being mindful of conversation history (if any).\n"
        "Synthesized Answer:"
    )
    messages_for_llm.append({"role": "user", "content": user_prompt_text})

    try:
        response = chat_client_instance.chat.completions.create(
            model=OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG,
            messages=messages_for_llm,
            temperature=0.1,
            max_tokens=1800
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        st.error(f"Tier 2 Snippet Synthesis Error (with history): {e}")
        error_fallback_ans = "Error during synthesis. Raw snippets:\n" + "\n---\n".join(
            [f"From {os.path.basename(s['source_doc'])}: {s['answer']}" for s in answer_snippets_with_sources]
        )
        return error_fallback_ans[:3000]

def select_best_tier2_answer(
        question_for_synthesis: str,
        tier2_positive_snippets: list[dict],
        chat_client_instance,
        conversation_history: list[dict] = None
    ) -> tuple[str | None, str | None]:
    if not tier2_positive_snippets:
        return None, None

    if len(tier2_positive_snippets) == 1:
        single_snippet = tier2_positive_snippets[0]
        return single_snippet["answer"], os.path.basename(single_snippet["source_doc"])

    synthesized_answer = synthesize_answers_from_snippets(
        question_for_synthesis,
        tier2_positive_snippets,
        chat_client_instance,
        conversation_history=conversation_history
    )

    if synthesized_answer and not is_answer_unsatisfactory(synthesized_answer):
        source_docs_involved = sorted(list(set(os.path.basename(s["source_doc"]) for s in tier2_positive_snippets)))
        source_ref_str = f"Synthesized from {len(source_docs_involved)} document(s): {', '.join(source_docs_involved)}"
        if len(source_ref_str) > 150:
            source_ref_str = f"Synthesized from {len(source_docs_involved)} document(s), including {source_docs_involved[0]}..."
        return synthesized_answer, source_ref_str
    else:
        return None, None

def process_uploaded_file_for_ingestion_and_validation(
    uploaded_file_object,
    ai_question_text: str,
    user_answer_text: str,
    kt_session_id: str,
    employee_name_val: str, # This will be APP_KT_USER_NAME
    uploader_identifier: str = "UserKTInterviewUpload"
) -> tuple[str | None, list[dict], str | None]:
    if not uploaded_file_object:
        return None, [], "No file object provided."

    file_name = uploaded_file_object.name
    file_bytes = uploaded_file_object.getvalue()
    file_extension = os.path.splitext(file_name)[-1].lower()
    extracted_text = ""
    st.info(f"Processing uploaded file for KT: {file_name}")

    try:
        if file_extension == ".pptx": extracted_text = extract_pptx_text(file_bytes)
        elif file_extension == ".docx": extracted_text = extract_docx_text(file_bytes)
        elif file_extension == ".doc": extracted_text = extract_doc_text(file_bytes, file_name)
        elif file_extension == ".pdf":
            pdf_parts = extract_pdf_text_and_images(file_bytes, file_name)
            extracted_text = "\n\n".join([p['content'] for p in pdf_parts if p['type'] == 'text' and p.get('content') and "Error:" not in p.get('content')])
        elif file_extension in [".xlsx", ".xls"]: extracted_text = extract_excel_text(file_bytes, file_extension)
        elif file_extension in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm', '.log', '.py', '.js', '.sql', ''] or \
             (not file_extension and file_name.lower() in KNOWN_EXTENSIONLESS_TEXT_FILES) :
            extracted_text = file_bytes.decode('utf-8', errors='ignore')
            if file_extension in ['.html', '.htm']:
                text_for_llm = re.sub(r'<style(?:\s[^>]*)?>.*?</style>', '', extracted_text, flags=re.DOTALL | re.IGNORECASE)
                text_for_llm = re.sub(r'<script(?:\s[^>]*)?>.*?</script>', '', text_for_llm, flags=re.DOTALL | re.IGNORECASE)
                text_for_llm = re.sub(r'<[^>]+>', ' ', text_for_llm)
                extracted_text = re.sub(r'\s+', ' ', text_for_llm).strip()
            extracted_text = remove_email_signature(extracted_text)
        else:
            return None, [], f"Unsupported file type for KT document: {file_extension}"

        if "Error:" in extracted_text or not extracted_text.strip():
            return None, [], f"Failed to extract text or empty content from {file_name}. Extractor output: {extracted_text[:100]}"

        chunks = semantic_chunk(extracted_text, max_chunk_len=1500, overlap=200)
        if not chunks:
            return extracted_text, [], f"No text chunks generated from {file_name} after processing."

        documents_for_search = []
        current_time_utc = datetime.now(timezone.utc)
        formatted_time = format_date_for_azure_search(current_time_utc)

        for i, chunk_content in enumerate(chunks):
            chunk_id_base = f"userupload_{kt_session_id}_{file_name}"
            chunk_id = get_stable_id(chunk_id_base, i)
            try:
                embedding_response = embed_client.embeddings.create(
                    input=[chunk_content], model=OPENAI_EMBED_LLM_DEPLOYMENT_NAME_INGEST
                )
                if not embedding_response.data or not embedding_response.data[0].embedding:
                    st.warning(f"Could not generate embedding for chunk {i} of {file_name}. Skipping chunk.")
                    continue
                vector = embedding_response.data[0].embedding
            except Exception as e_embed:
                st.warning(f"Embedding error for chunk {i} of {file_name}: {e_embed}. Skipping chunk.")
                continue

            doc = {
                "id": chunk_id, "chunkText": chunk_content, "vector": vector,
                "sourceFileName": file_name,
                "sourceFilePath": f"{KT_UPLOAD_PATH_PREFIX}/{kt_session_id}/{file_name}",
                "extension": file_extension,
                "documentType": "UserKTSupportDocument",
                "uploaded_by": uploader_identifier,
                "upload_date": formatted_time,
                "timestamp": formatted_time,
                "source": "KTInterviewUpload",
                "chunkIndex": i,
                "employeeName": employee_name_val, # This should be APP_KT_USER_NAME
                "kt_session_id": kt_session_id,
                "kt_session_question": ai_question_text[:1000] if ai_question_text else None,
                "kt_session_answer": user_answer_text[:1000] if user_answer_text else None,
            }
            documents_for_search.append(doc)
        st.success(f"Successfully processed and chunked {file_name} for KT session into {len(documents_for_search)} vectorizable chunks.")
        return extracted_text, documents_for_search, None
    except Exception as e:
        st.error(f"Error processing uploaded KT file {file_name}: {e}")
        traceback.print_exc()
        return None, [], f"Fatal error processing KT file {file_name}: {e}"

def validate_answer_with_llm(
    ai_question: str, user_answer: str, document_context: str, source_description: str
) -> tuple[str, float | None]:
    system_prompt = (
        "You are an expert evaluator. Your task is to assess if the 'Provided Document Excerpt' "
        "supports the 'User's Answer' in response to the 'Original AI Question'. Be objective and concise."
        "First, state if the excerpt supports, partially supports, contradicts, or is irrelevant to the user's answer. "
        "Then, briefly explain why. Finally, on a new line, provide a numerical score from 0.0 (no support) to 1.0 (full direct support), "
        "like: SCORE: 0.8"
        "If the excerpt is irrelevant, the score must be 0.0."
        "Consider if the user's answer directly addresses the AI question based on the excerpt."
    )
    user_prompt = (
        f"Original AI Question: \"{ai_question}\"\n\n"
        f"User's Answer: \"{user_answer}\"\n\n"
        f"Provided Document Excerpt (from {source_description}):\n"
        f"```text\n{document_context[:7000]}\n```\n\n"
        "Evaluation (Support statement, explanation, and SCORE: X.X):"
    )
    try:
        response = chat_client.chat.completions.create(
            model=OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG,
            messages=[ {"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt} ],
            temperature=0.1,
            max_tokens=500
        )
        content = response.choices[0].message.content.strip()

        score_match = re.search(r"SCORE:\s*([0-1](?:\.\d+)?)", content, re.IGNORECASE)
        parsed_score = float(score_match.group(1)) if score_match else None

        summary_text = re.sub(r"SCORE:\s*([0-1](?:\.\d+)?)", "", content, flags=re.IGNORECASE).strip()
        if not summary_text and parsed_score is not None:
            summary_text = f"Validation score is {parsed_score:.2f}."
        elif not summary_text and parsed_score is None:
            summary_text = "LLM evaluation did not provide a clear summary or score."

        return summary_text, parsed_score
    except Exception as e:
        st.error(f"LLM validation error: {e}")
        traceback.print_exc()
        return f"Error during LLM validation: {e}", None

def find_relevant_indexed_chunks(
    query_text: str, filter_expression: str = None, top_k: int = 1
) -> list[dict]:
    if not query_text: return []
    try:
        query_vector_response = embed_client.embeddings.create(
            input=[query_text], model=OPENAI_EMBED_LLM_DEPLOYMENT_NAME_INGEST
        )
        if not query_vector_response.data or not query_vector_response.data[0].embedding:
            st.error("Validation RAG: Embedding Error - No vector received.")
            return []
        query_vector = query_vector_response.data[0].embedding
    except Exception as e:
        st.error(f"Validation RAG: Embedding Error: {e}")
        return []

    vector_query = VectorizedQuery(vector=query_vector, k_nearest_neighbors=top_k, fields="vector", exhaustive=True)
    select_fields = ["id", "chunkText", "sourceFileName", "documentType", "kt_session_question", "kt_session_answer", "projectName"]

    try:
        results = search_client.search(
            search_text=query_text,
            vector_queries=[vector_query],
            filter=filter_expression,
            select=select_fields,
            top=top_k
        )
        docs = []
        for res in results:
            doc = {field: res.get(field) for field in select_fields if field in res}
            doc["@search.score"] = res.get("@search.score", 0.0)
            docs.append(doc)

        docs.sort(key=lambda x: x.get("@search.score", 0.0), reverse=True)
        return docs
    except Exception as e:
        st.error(f"Validation RAG: Search Error for query '{query_text[:50]}...': {e}")
        return []

def perform_tier2_validation_deep_dive(
    ai_question: str, user_answer: str,
    files_to_search: list,
    chat_client_instance, blob_service_client_instance, adls_container: str
) -> tuple[str | None, str | None, float | None]:
    st.write(f"  Tier 2 Validation (KT): Analyzing {len(files_to_search)} existing files for user's answer...")
    best_overall_snippet_text = None
    best_snippet_source_filename = None
    highest_score = -1.0

    system_prompt_tier2_chunk_eval = (
        "You are a meticulous AI information evaluator. Your task is to determine if the 'Text Snippet' "
        "contains information that supports, contradicts, or is irrelevant to the 'User's Answer' "
        "which was given in response to the 'Original AI Question'.\n"
        "1. If the snippet supports the user's answer (fully or partially), extract the supporting text verbatim or as a close paraphrase. Then, on a new line, provide a numerical score from 0.5 (partial/weak support) to 1.0 (strong/direct support), like: SCORE: 0.8.\n"
        "2. If the snippet contradicts the user's answer, state that and explain briefly. Then, on a new line, provide SCORE: 0.0.\n"
        "3. If the snippet is irrelevant to the user's answer regarding the AI question, respond ONLY with the exact phrase: " + DESIGNATED_NO_ANSWER_PHRASE + " and on a new line, SCORE: 0.0.\n"
        "Focus ONLY on the relationship between the snippet and the USER'S ANSWER for the given AI Question."
    )

    processed_file_count = 0
    for blob_item_data in files_to_search:
        blob_name = blob_item_data["name"]
        base_filename = os.path.basename(blob_name)
        blob_size = blob_item_data.get("size", 0)

        if blob_size == 0: continue

        try:
            blob_data_bytes = read_blob_to_memory(blob_service_client_instance, adls_container, blob_name)
            file_extension = os.path.splitext(blob_name)[-1].lower()
            full_text_content_for_doc = ""

            extractor_map = { '.pptx': extract_pptx_text, '.docx': extract_docx_text }
            if file_extension in extractor_map: full_text_content_for_doc = extractor_map[file_extension](blob_data_bytes)
            elif file_extension == '.doc': full_text_content_for_doc = extract_doc_text(blob_data_bytes, blob_name)
            elif file_extension == '.pdf':
                pdf_parts = extract_pdf_text_and_images(blob_data_bytes, blob_name)
                full_text_content_for_doc = "\n\n".join([p['content'] for p in pdf_parts if p['type'] == 'text' and p.get('content') and "Error:" not in p.get('content')])
            elif file_extension in ['.xlsx', '.xls']: full_text_content_for_doc = extract_excel_text(blob_data_bytes, file_extension)
            elif file_extension in ['.txt', '.md', '.csv', '.json', '.xml', '.html', '.htm', '.log', ''] or \
                 (not file_extension and base_filename.lower() in KNOWN_EXTENSIONLESS_TEXT_FILES):
                decoded_text = blob_data_bytes.decode('utf-8', errors='ignore')
                if file_extension in ['.html', '.htm']:
                    text_llm = re.sub(r'<style[^>]*>.*?</style>|<script[^>]*>.*?</script>', '', decoded_text, flags=re.S|re.I)
                    text_llm = re.sub(r'<[^>]+>', ' ', text_llm)
                    full_text_content_for_doc = re.sub(r'\s+', ' ', text_llm).strip()
                else: full_text_content_for_doc = decoded_text
                full_text_content_for_doc = remove_email_signature(full_text_content_for_doc)
            else: continue

            if "Error:" in full_text_content_for_doc or not full_text_content_for_doc.strip(): continue

            tier2_semantic_chunks = semantic_chunk(full_text_content_for_doc, TIER2_CHUNK_SIZE, TIER2_CHUNK_OVERLAP)
            if not tier2_semantic_chunks: continue

            for chunk_idx, chunk_text in enumerate(tier2_semantic_chunks[:MAX_CHUNKS_PER_DOC_TIER2]):
                if not chunk_text.strip(): continue

                user_prompt_chunk_eval = (
                    f"Original AI Question: \"{ai_question}\"\n\n"
                    f"User's Answer: \"{user_answer}\"\n\n"
                    f"Text Snippet (from document '{base_filename}', chunk {chunk_idx+1}):\n"
                    f"```text\n{chunk_text}\n```\n\n"
                    "Evaluation (Support/Contradict/Irrelevant statement, explanation, and SCORE: X.X):"
                )

                response_chunk_eval = chat_client_instance.chat.completions.create(
                    model=OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG,
                    messages=[ {"role": "system", "content": system_prompt_tier2_chunk_eval}, {"role": "user", "content": user_prompt_chunk_eval} ],
                    temperature=0.0, max_tokens=400
                )
                eval_content = response_chunk_eval.choices[0].message.content.strip()

                if DESIGNATED_NO_ANSWER_PHRASE not in eval_content:
                    score_match = re.search(r"SCORE:\s*([0-1](?:\.\d+)?)", eval_content, re.IGNORECASE)
                    current_score = float(score_match.group(1)) if score_match else 0.0

                    if current_score > 0.0 and current_score > highest_score:
                        highest_score = current_score
                        best_overall_snippet_text = re.sub(r"SCORE:\s*([0-1](?:\.\d+)?)", "", eval_content, flags=re.IGNORECASE).strip()
                        best_snippet_source_filename = base_filename
        except Exception as e_tier2_file:
            st.warning(f"Tier 2 Validation: Error processing file {blob_name} for KT: {e_tier2_file}")

        processed_file_count += 1
        if processed_file_count % 5 == 0:
            st.caption(f"    Tier 2 Validation (KT): Processed {processed_file_count}/{len(files_to_search)} files...")

    if best_overall_snippet_text and highest_score > 0:
        st.success(f"Tier 2 Validation (KT): Found supporting information in '{best_snippet_source_filename}' with score {highest_score:.2f}.")
        return best_overall_snippet_text, best_snippet_source_filename, highest_score
    else:
        st.info("Tier 2 Validation (KT): No strong supporting information found in existing raw documents.")
        return None, None, 0.0
# --- END RAG, Validation, etc. ---


# --- Streamlit App UI and Main Orchestration Logic ---
st.title("🤖 CogniLink: AI Knowledge Transfer & Validation Pipeline")

# --- Session State Initialization ---
if 'questions' not in st.session_state: st.session_state.questions = []
if 'answers' not in st.session_state: st.session_state.answers = {}
if 'pipeline_started' not in st.session_state: st.session_state.pipeline_started = False
if 'tier2_triggered_for' not in st.session_state: st.session_state.tier2_triggered_for = set()
if 'qa_summary_text' not in st.session_state: st.session_state.qa_summary_text = ""
if 'qa_summary_audio_path' not in st.session_state: st.session_state.qa_summary_audio_path = ""

if 'chat_mode' not in st.session_state: st.session_state.chat_mode = False
if 'chat_history' not in st.session_state: st.session_state.chat_history = []

if 'kt_interview_active' not in st.session_state: st.session_state.kt_interview_active = False
if 'kt_interview_conversation_history' not in st.session_state: st.session_state.kt_interview_conversation_history = []
if 'current_kt_ai_question_details' not in st.session_state: st.session_state.current_kt_ai_question_details = None
if 'kt_interview_question_count' not in st.session_state: st.session_state.kt_interview_question_count = 0
if 'kt_session_id' not in st.session_state: st.session_state.kt_session_id = None
if 'kt_interview_results' not in st.session_state: st.session_state.kt_interview_results = []
if 'max_kt_questions_user_set' not in st.session_state: st.session_state.max_kt_questions_user_set = 15
if 'last_validation_display' not in st.session_state: st.session_state.last_validation_display = None # For KT interview validation UI


if 'all_blobs_for_kt_target_user' not in st.session_state: st.session_state.all_blobs_for_kt_target_user = []
if 'debug_llm_messages_kt_qgen' not in st.session_state: st.session_state.debug_llm_messages_kt_qgen = []
if 'debug_kt_qgen_error_response' not in st.session_state: st.session_state.debug_kt_qgen_error_response = None


# --- Global Variables for App Context (Derived from Environment Variables) ---
PIPELINE_SUBJECT_NAME_FOR_QGEN = f"{APP_KT_USER_NAME} ({APP_KT_USER_EMAIL})"
DESCRIPTION_FOR_Q_GEN = (
    f"Key knowledge, projects, specific contributions, responsibilities, technical skills, KPIs, and insights "
    f"derived from SharePoint documents or KT sessions related to {PIPELINE_SUBJECT_NAME_FOR_QGEN}."
)

# --- Initial Data Loading ---
@st.cache_data(ttl=3600)
def load_kt_target_user_blob_list_cached():
    st.info(f"Attempting to load file list for user '{APP_TARGET_USER_EMAIL_FOR_BLOBS}' from their source data folder ('{APP_TARGET_SOURCE_FOLDER_FOR_BLOBS}') for Tier 2 deep dives...")
    blobs_found_list = []
    if blob_service_client and ADLS_CONTAINER_NAME:
        try:
            container_client = blob_service_client.get_container_client(ADLS_CONTAINER_NAME)
            # Ensure prefix ends with '/' if APP_TARGET_SOURCE_FOLDER_FOR_BLOBS is not empty,
            # and handles case where it might already include user email or be a direct path.
            prefix_to_scan = APP_TARGET_SOURCE_FOLDER_FOR_BLOBS
            if APP_TARGET_SOURCE_FOLDER_FOR_BLOBS and not APP_TARGET_SOURCE_FOLDER_FOR_BLOBS.endswith('/'):
                prefix_to_scan += '/'
            # If the folder structure convention implies user email is part of it:
            # prefix_to_scan = f"{APP_TARGET_SOURCE_FOLDER_FOR_BLOBS}/{APP_TARGET_USER_EMAIL_FOR_BLOBS}/"
            # The current env var APP_TARGET_SOURCE_FOLDER_FOR_BLOBS is expected to be the full path prefix needed.

            blob_list_iterable = container_client.list_blobs(name_starts_with=prefix_to_scan)
            for blob_properties in blob_list_iterable:
                if not blob_properties.name.endswith('/') and blob_properties.size and blob_properties.size > 0:
                    blobs_found_list.append({
                        "name": blob_properties.name,
                        "size": blob_properties.size,
                        "last_modified": blob_properties.last_modified
                    })
            if blobs_found_list:
                st.success(f"File list loaded: {len(blobs_found_list)} relevant files found under prefix '{prefix_to_scan}'. These can be used for Tier 2 deep dives.")
            else:
                st.warning(f"No files found with prefix '{prefix_to_scan}'. Tier 2 deep dive for existing documents might not be effective if this path is incorrect or empty.")
            return blobs_found_list
        except Exception as e_blob_list:
            st.error(f"Error listing blobs from ADLS container '{ADLS_CONTAINER_NAME}' with prefix '{prefix_to_scan}': {e_blob_list}")
            traceback.print_exc()
            return []
    else:
        st.error("Azure Blob Service client or ADLS Container Name not configured. Cannot list files for Tier 2.")
        return []

if not st.session_state.all_blobs_for_kt_target_user:
    st.session_state.all_blobs_for_kt_target_user = load_kt_target_user_blob_list_cached()


# --- Sidebar for Mode Selection ---
with st.sidebar:
    st.header("Operation Mode")
    pipeline_choice = st.radio(
        "Choose mode:",
        ("Automated Q&A Pipeline", "Structured KT Interview (CogniLink)", "Interactive Chat"),
        key="pipeline_choice_radio",
        horizontal=False,
        label_visibility="collapsed"
    )

    if 'last_pipeline_choice' not in st.session_state:
        st.session_state.last_pipeline_choice = pipeline_choice

    if st.session_state.last_pipeline_choice != pipeline_choice:
        st.session_state.pipeline_started = False
        st.session_state.kt_interview_active = False
        st.session_state.chat_mode = False
        st.session_state.last_validation_display = None
        st.session_state.last_pipeline_choice = pipeline_choice
        st.rerun()

    if pipeline_choice == "Automated Q&A Pipeline":
        st.session_state.pipeline_started = True
        st.session_state.kt_interview_active = False
        st.session_state.chat_mode = False
        st.subheader("Automated Q&A Controls")
        num_q_gen_input_auto = st.number_input("Number of AI Questions to Generate:", min_value=1, max_value=10, value=3, step=1, key="num_q_ctrl_auto")
        if st.button("🚀 Start Automated Knowledge Pipeline", key="start_btn_auto_pipeline"):
            st.session_state.questions = []
            st.session_state.answers = {}
            st.session_state.tier2_triggered_for = set()
            st.session_state.qa_summary_text = ""
            st.session_state.qa_summary_audio_path = ""
            st.session_state.last_validation_display = None

            if not st.session_state.all_blobs_for_kt_target_user: # Changed from all_blobs_for_gaurav
                st.session_state.all_blobs_for_kt_target_user = load_kt_target_user_blob_list_cached()

            with st.spinner("Generating AI-powered questions for automated pipeline..."):
                st.session_state.questions = generate_ai_questions_from_indexed_data(
                    focus_entity_name=PIPELINE_SUBJECT_NAME_FOR_QGEN,
                    entity_type_description=DESCRIPTION_FOR_Q_GEN,
                    num_questions=num_q_gen_input_auto,
                    use_rag_summary_context=True,
                    search_filter_for_summary=KT_TARGET_RELEVANT_CONTENT_FILTER # Changed
                )
            if not st.session_state.questions:
                st.error("Failed to generate questions for Automated Pipeline.")
            else:
                st.success(f"Generated {len(st.session_state.questions)} questions. Processing automated Q&A...")
                progress_bar = st.progress(0)
                for q_idx, current_question_text in enumerate(st.session_state.questions):
                    question_label = f"Q{q_idx+1}/{len(st.session_state.questions)}"
                    st.markdown(f"--- \n**Processing {question_label}:** *{current_question_text[:100]}...*")

                    with st.spinner(f"{question_label}: Tier 1 RAG answering..."):
                        tier1_answer, found_context_t1, retrieved_docs_t1 = answer_question_with_rag(
                            current_question_text,
                            KT_TARGET_RELEVANT_CONTENT_FILTER, # Changed
                            top_k_contexts=5
                        )

                    final_answer_text = tier1_answer
                    answer_tier_description = "Tier 1 (RAG from Search Index)"
                    answer_source_docs = ", ".join(list(set(os.path.basename(doc.get("sourceFileName", "Unknown")) for doc in retrieved_docs_t1 if doc.get("sourceFileName")))) if retrieved_docs_t1 else "N/A"

                    if not found_context_t1 or is_answer_unsatisfactory(tier1_answer):
                        st.warning(f"{question_label}: Tier 1 unsatisfactory or no context. Initiating Tier 2 Deep Dive...")
                        st.session_state.tier2_triggered_for.add(q_idx)
                        if not st.session_state.all_blobs_for_kt_target_user: # Changed
                            st.error(f"{question_label}: Tier 2 files unavailable. Sticking with Tier 1 result.")
                            final_answer_text = tier1_answer if found_context_t1 else "Information not found (Tier 1 & Tier 2 files unavailable)."
                            answer_tier_description = "Tier 1 (No files for Tier 2)" if found_context_t1 else "Undocumented (Tier 2 files unavailable)"
                        else:
                            with st.spinner(f"{question_label}: Tier 2 Deep Dive in progress... (analyzing raw files)"):
                                tier2_positive_snippets = perform_tier2_deep_dive_for_one_question(
                                    current_question_text, st.session_state.all_blobs_for_kt_target_user, # Changed
                                    chat_client, blob_service_client, ADLS_CONTAINER_NAME
                                )
                            if tier2_positive_snippets:
                                with st.spinner(f"{question_label}: Synthesizing Tier 2 answer..."):
                                    best_tier2_answer, tier2_source_ref = select_best_tier2_answer(
                                        current_question_text, tier2_positive_snippets, chat_client
                                    )
                                if best_tier2_answer:
                                    final_answer_text = best_tier2_answer
                                    answer_tier_description = "Tier 2 (Deep Dive from Raw Files)"
                                    answer_source_docs = tier2_source_ref
                                else:
                                    final_answer_text = tier1_answer if found_context_t1 and not is_answer_unsatisfactory(tier1_answer) else "Information remains elusive after Tier 1 & Tier 2."
                                    answer_tier_description = "Tier 1 (Tier 2 unsatisfactory)" if found_context_t1 and not is_answer_unsatisfactory(tier1_answer) else "Undocumented (After Tier 1 & Tier 2)"
                            else:
                                final_answer_text = tier1_answer if found_context_t1 and not is_answer_unsatisfactory(tier1_answer) else "No specific information found after Tier 1 & Tier 2 search."
                                answer_tier_description = "Tier 1 (No Tier 2 snippets)" if found_context_t1 and not is_answer_unsatisfactory(tier1_answer) else "Undocumented (No Tier 2 snippets)"

                    st.session_state.answers[q_idx] = {
                        "question": current_question_text, "answer": final_answer_text,
                        "tier": answer_tier_description, "sources": answer_source_docs if answer_source_docs else "N/A"
                    }

                    combined_q_and_a_text_for_ingestion = f"Question: {current_question_text}\nAnswer: {final_answer_text}"
                    try:
                        with st.spinner(f"{question_label}: Storing Q&A to knowledge base..."):
                            embedding_response_qna = embed_client.embeddings.create(input=[combined_q_and_a_text_for_ingestion], model=OPENAI_EMBED_LLM_DEPLOYMENT_NAME_INGEST)
                            if embedding_response_qna.data and embedding_response_qna.data[0].embedding:
                                qna_doc_id = sanitize_id_for_search(f"ktsession_autoqna_{APP_KT_USER_EMAIL}_{int(time.time())}_q{q_idx}") # Changed
                                current_time_iso = format_date_for_azure_search(datetime.now(timezone.utc))
                                new_qna_document_for_search = {
                                    "id": qna_doc_id,
                                    "vector": embedding_response_qna.data[0].embedding,
                                    "kt_session_question": current_question_text,
                                    "kt_session_answer": final_answer_text,
                                    "chunkText": combined_q_and_a_text_for_ingestion,
                                    "documentType": "KnowledgeTransferSessionQA",
                                    "employeeName": APP_KT_USER_EMAIL, # Changed
                                    "sourceFileName": f"AutoQ&A_Session_{APP_KT_USER_NAME.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json", # Changed
                                    "uploaded_by": "AISystemKnowledgePipeline",
                                    "upload_date": current_time_iso,
                                    "timestamp": current_time_iso,
                                    "source": "InternalQASynthesis"
                                }
                                search_client.upload_documents(documents=[new_qna_document_for_search])
                                st.success(f"{question_label}: Q&A pair from automated pipeline stored in knowledge base.")
                            else: st.error(f"{question_label}: Failed to embed Q&A for storage.")
                    except Exception as e_embed_store_qna:
                        st.error(f"{question_label}: Error storing Q&A to knowledge base: {e_embed_store_qna}")

                    progress_bar.progress((q_idx + 1) / len(st.session_state.questions))
                st.balloons()
                st.success("Automated Knowledge Pipeline Completed!")
            st.rerun()

        if st.session_state.questions and st.session_state.answers:
            st.markdown("---")
            if st.button("🗑 Reset Automated Pipeline Results", key="reset_btn_auto_pipeline"):
                st.session_state.questions = []
                st.session_state.answers = {}
                st.session_state.tier2_triggered_for = set()
                st.session_state.qa_summary_text = ""
                st.session_state.qa_summary_audio_path = ""
                st.session_state.last_validation_display = None
                st.warning("Automated pipeline results cleared.")
                time.sleep(1); st.rerun()

            st.markdown("---")
            if st.button("🔊 Generate Q&A Summary (Text + Audio)", key="summary_btn_auto", disabled=not TTS_AVAILABLE):
                if not st.session_state.answers: st.warning("No Q&A to summarize.")
                else:
                    summary_prompt_parts = [
                        f"Generate a concise executive summary of the key information extracted about {PIPELINE_SUBJECT_NAME_FOR_QGEN} "
                        "from the following Question and Answer pairs. Focus on projects, achievements, skills, and critical knowledge. "
                        "Be factual and directly cite the sources mentioned if important.\n\n"
                    ]
                    for idx, qa_data in st.session_state.answers.items():
                        summary_prompt_parts.append(f"Q{idx+1}: {qa_data['question']}\nA{idx+1} (Tier: {qa_data['tier']}, Source(s): {qa_data.get('sources','N/A')}):\n{qa_data['answer']}\n---\n")
                    full_summary_prompt = "".join(summary_prompt_parts)

                    with st.spinner("Generating text summary..."):
                        try:
                            summary_response = chat_client.chat.completions.create(
                                model=OPENAI_CHAT_LLM_DEPLOYMENT_NAME_RAG,
                                messages=[{"role": "system", "content": "You are an expert summarizer of Q&A sessions."},
                                          {"role": "user", "content": full_summary_prompt}],
                                temperature=0.2,
                                max_tokens=1500
                            )
                            st.session_state.qa_summary_text = summary_response.choices[0].message.content.strip()
                            st.success("Text summary generated.")
                        except Exception as e_summary_gen:
                            st.error(f"Summary generation error: {e_summary_gen}")
                            st.session_state.qa_summary_text = f"Error generating summary: {e_summary_gen}"

                    if st.session_state.qa_summary_text and "Error:" not in st.session_state.qa_summary_text and TTS_AVAILABLE:
                        with st.spinner("Generating audio summary... (This may take a moment)"):
                            try:
                                tts_instance = TTS(model_name=TTS_MODEL_NAME, progress_bar=False) # Use env var for model name
                                temp_audio_file_path = os.path.join(tempfile.gettempdir(), f"summary_auto_{APP_KT_USER_NAME.replace(' ', '_')}_{int(time.time())}.wav") # Changed
                                text_for_audio = st.session_state.qa_summary_text
                                if len(text_for_audio.split()) > 300:
                                    text_for_audio = " ".join(text_for_audio.split()[:300]) + "... Summary truncated for audio."
                                    st.caption("Audio summary may be truncated due to length.")

                                tts_instance.tts_to_file(text=text_for_audio, file_path=temp_audio_file_path)
                                st.session_state.qa_summary_audio_path = temp_audio_file_path
                                st.success("Audio summary generated.")
                            except Exception as e_tts:
                                st.error(f"TTS error: {e_tts}. Coqui TTS might need specific setup or model download (Model: {TTS_MODEL_NAME}).")
                                st.session_state.qa_summary_audio_path = ""
                st.rerun()
            elif not TTS_AVAILABLE:
                 st.caption("Audio summary disabled as TTS library is not available.")


    elif pipeline_choice == "Structured KT Interview (CogniLink)":
        if 'max_kt_questions_user_set' not in st.session_state:
            st.session_state.max_kt_questions_user_set = 15

        st.session_state.kt_interview_active = True
        st.session_state.pipeline_started = False
        st.session_state.chat_mode = False
        st.subheader(f"CogniLink KT Interview with {APP_KT_USER_NAME}") # Changed

        st.session_state.max_kt_questions_user_set = st.number_input(
            "Max Interview Turns (AI Questions):",
            min_value=5,
            max_value=50,
            value=st.session_state.max_kt_questions_user_set,
            step=1,
            key="max_kt_questions_input",
            help="Sets the maximum number of questions CogniLink will ask. Each AI utterance, including sub-questions, counts as one turn."
        )

        if st.button("🚀 Start New CogniLink KT Interview", key="start_btn_kt_interview"):
            st.session_state.kt_interview_active = True
            st.session_state.kt_interview_conversation_history = []
            st.session_state.current_kt_ai_question_details = None
            st.session_state.kt_interview_question_count = 0
            st.session_state.kt_interview_results = []
            st.session_state.kt_session_id = f"ktsession_cognilink_{APP_KT_USER_EMAIL.split('@')[0]}_{int(time.time())}" # Changed
            st.session_state.debug_llm_messages_kt_qgen = []
            st.session_state.debug_kt_qgen_error_response = None
            st.session_state.last_validation_display = None

            with st.spinner("CogniLink is preparing the first question..."):
                st.session_state.current_kt_ai_question_details = generate_single_kt_interview_question(
                    employee_name=APP_KT_USER_NAME, # Changed
                    employee_email=APP_KT_USER_EMAIL, # Changed
                    conversation_history=[],
                )
            if st.session_state.current_kt_ai_question_details and "error" not in st.session_state.current_kt_ai_question_details.get("question","").lower():
                # Add the first AI question to history (as it's now asked)
                st.session_state.kt_interview_conversation_history.append({
                    "role": "assistant",
                    "content": json.dumps(st.session_state.current_kt_ai_question_details)
                })
                st.session_state.kt_interview_question_count = 1 # First question asked
            else:
                st.session_state.kt_interview_active = False
                st.error("Could not start KT interview: CogniLink had an issue generating the first question.")
            st.rerun()

        if st.session_state.kt_interview_active and st.session_state.current_kt_ai_question_details:
            if st.button("⏹️ End and Reset CogniLink Interview", key="reset_btn_kt_interview"):
                st.session_state.kt_interview_active = False
                st.session_state.kt_interview_conversation_history = []
                st.session_state.current_kt_ai_question_details = None
                st.session_state.kt_interview_question_count = 0
                st.session_state.kt_interview_results = []
                st.session_state.kt_session_id = None
                st.session_state.last_validation_display = None
                st.warning("CogniLink KT Interview session reset.")
                time.sleep(1); st.rerun()

    elif pipeline_choice == "Interactive Chat":
        st.session_state.chat_mode = True
        st.session_state.pipeline_started = False
        st.session_state.kt_interview_active = False
        st.session_state.last_validation_display = None
        st.subheader(f"Interactive Chat with {APP_KT_USER_NAME}'s Knowledge Base") # Changed
        st.write("Ask questions directly. The AI will use existing indexed knowledge, including previous KT sessions and uploaded documents.")

# --- Main Content Area ---
if st.session_state.pipeline_started:
    st.header("Automated Q&A Pipeline Results")
    if st.session_state.answers:
        for q_idx_display, qa_info in st.session_state.answers.items():
            with st.expander(f"Q{q_idx_display+1}: {qa_info['question']}", expanded=True):
                st.markdown(f"**Answer (Source Tier: {qa_info['tier']}):**")
                if "Undocumented" in qa_info["tier"] or is_answer_unsatisfactory(qa_info["answer"]):
                    st.warning(f"{qa_info['answer']}")
                else:
                    st.success(f"{qa_info['answer']}")
                st.caption(f"Retrieved Source(s): {qa_info.get('sources', 'N/A')}")
                if q_idx_display in st.session_state.tier2_triggered_for:
                    st.info("ℹ️ Tier 2 Deep Dive (raw file scan) was performed for this question.")

        if st.session_state.qa_summary_text:
            st.markdown("---"); st.header("Executive Summary of Automated Q&A")
            st.markdown(st.session_state.qa_summary_text)
            if st.session_state.qa_summary_audio_path and os.path.exists(st.session_state.qa_summary_audio_path) and TTS_AVAILABLE:
                try:
                    with open(st.session_state.qa_summary_audio_path, "rb") as audio_f:
                        st.audio(audio_f.read(), format="audio/wav")
                except Exception as e_audio_play: st.error(f"Could not play audio summary: {e_audio_play}")
            elif not TTS_AVAILABLE:
                st.caption("Audio summary generation is disabled.")

    elif st.session_state.questions:
         st.info("Automated pipeline started. Processing questions...")
    else:
        st.info("Start the 'Automated Q&A Pipeline' from the sidebar to generate and answer questions automatically.")


elif st.session_state.kt_interview_active:
    MAX_KT_INTERVIEW_QUESTIONS_TO_ASK = st.session_state.max_kt_questions_user_set

    # Display last validation result if available
    if st.session_state.last_validation_display:
        val_info = st.session_state.last_validation_display
        if val_info.get('score', -1) >= 0.7:
            st.success(f"System Validation: Score {val_info['score']:.2f} (Source: {val_info['source']}). Details: {val_info['details']}")
        elif val_info.get('score', -1) >= 0.4:
            st.info(f"System Validation: Score {val_info['score']:.2f} (Source: {val_info['source']}). Details: {val_info['details']}")
        else:
            st.warning(f"System Validation: Score {val_info.get('score', 0.0):.2f} (Source: {val_info['source']}). Details: {val_info['details']}")
        st.session_state.last_validation_display = None # Clear after displaying once

    if st.session_state.current_kt_ai_question_details and \
       st.session_state.kt_interview_question_count <= MAX_KT_INTERVIEW_QUESTIONS_TO_ASK and \
       "error" not in st.session_state.current_kt_ai_question_details.get("question","").lower() and \
       "[END_OF_INTERVIEW]" not in st.session_state.current_kt_ai_question_details.get("question",""):

        q_details = st.session_state.current_kt_ai_question_details
        q_count = st.session_state.kt_interview_question_count

        st.header(f"CogniLink Interview: Turn {q_count}/{MAX_KT_INTERVIEW_QUESTIONS_TO_ASK}")
        st.subheader(f"CogniLink: \"{q_details['question']}\"")

        with st.expander("View Suggested Response Structure (Guide)", expanded=False):
            st.json(q_details["response_structure"])

        current_q_uploaded_doc_name = None
        current_q_processed_chunks = []

        with st.form(key=f"kt_answer_form_q{q_count}"):
            user_answer_text_key = f"kt_user_ans_text_q{q_count}"
            user_answer_text_val = st.text_area(f"Your Answer (or type 'next' to move on):", height=150, key=user_answer_text_key)

            uploaded_support_doc_obj = st.file_uploader(
                "Upload a supporting document for this answer (optional):",
                type=['txt', 'pdf', 'docx', 'doc', 'pptx', 'xlsx', 'xls', 'md', 'json'],
                key=f"kt_user_upload_q{q_count}"
            )
            submit_kt_answer_button = st.form_submit_button("✅ Submit Answer & Continue Interview")

        if submit_kt_answer_button:
            final_validation_score = 0.0 # Default score
            final_validation_source = "N/A"
            validation_details_text = "No validation performed yet or answer too short."
            ingested_doc_msg = None


            if user_answer_text_val.strip().lower() == "next":
                st.session_state.kt_interview_conversation_history.append({
                    "role": "system",
                    "content": json.dumps({
                        "observation": "User has indicated they wish to move to the next main topic or question. Please proceed accordingly in the interview flow, concluding the current deep dive if appropriate."
                    })
                })
                st.session_state.kt_interview_conversation_history.append({
                    "role": "user",
                    "content": "next"
                })
                st.session_state.kt_interview_results.append({
                    "kt_session_id": st.session_state.kt_session_id,
                    "question_number": q_count,
                    "ai_question_details": dict(q_details),
                    "user_text_answer": "next (User initiated topic change)",
                    "final_validation_source": "User Action",
                    "final_validation_details": "User typed 'next'.",
                    "final_validation_score": 1.0,
                    "uploaded_document_name": None,
                    "ingestion_status_message": "N/A (User typed 'next')",
                    "raw_azure_search_score_if_rag": None,
                    "timestamp": format_date_for_azure_search(datetime.now(timezone.utc))
                })
                st.info("Moving to the next topic as requested...")

            elif user_answer_text_val.strip():
                current_ai_q_text = q_details['question']

                st.session_state.kt_interview_conversation_history.append({
                    "role": "user",
                    "content": user_answer_text_val
                })

                # Initialize validation variables
                doc_validation_score = 0.0
                doc_validation_details = "No document uploaded or processed."
                doc_validation_source_file = None
                raw_azure_search_score_from_rag = None # Initialize

                # 1. Validate against user-uploaded document
                if uploaded_support_doc_obj:
                    current_q_uploaded_doc_name = uploaded_support_doc_obj.name
                    doc_validation_source_file = current_q_uploaded_doc_name # Keep track of the source file
                    with st.spinner(f"Processing & validating with uploaded document: {current_q_uploaded_doc_name}..."):
                        doc_text_content, temp_chunks_for_ingest, err_msg = process_uploaded_file_for_ingestion_and_validation(
                            uploaded_support_doc_obj,
                            current_ai_q_text,
                            user_answer_text_val,
                            st.session_state.kt_session_id,
                            APP_KT_USER_NAME # Changed
                        )
                    if err_msg:
                        st.error(f"Error processing '{current_q_uploaded_doc_name}': {err_msg}")
                        doc_validation_details = f"Error processing document: {err_msg}"
                    elif doc_text_content:
                        current_q_processed_chunks = temp_chunks_for_ingest
                        val_summary_doc, val_score_doc_raw = validate_answer_with_llm(
                            current_ai_q_text, user_answer_text_val, doc_text_content, current_q_uploaded_doc_name
                        )
                        if val_score_doc_raw is not None:
                            doc_validation_score = val_score_doc_raw
                            doc_validation_details = val_summary_doc
                            st.success(f"Validation with '{current_q_uploaded_doc_name}' complete. LLM Score: {doc_validation_score:.2f}")
                        else:
                            doc_validation_details = "LLM validation failed for uploaded document."
                            st.warning(doc_validation_details)
                    else:
                        st.warning(f"Could not extract significant content from '{current_q_uploaded_doc_name}'.")
                        doc_validation_details = "No significant content extracted from document."

                # Tentative final validation based on doc or self-attestation
                final_validation_score = doc_validation_score
                final_validation_source = f"User Upload: {doc_validation_source_file}" if doc_validation_source_file else "User Provided Text (No Doc)"
                validation_details_text = doc_validation_details if doc_validation_source_file else "User self-attested answer."
                if not doc_validation_source_file: # If no document was uploaded or processed successfully
                    final_validation_score = 0.5 # Base score for self-attested if no doc

                # 2. Validate against RAG and capture raw Azure Search score
                rag_llm_validation_score = 0.0       # LLM's validation score for the RAG chunk's content

                with st.spinner("Validating with existing knowledge base (RAG)..."):
                    rag_query = f"Regarding the question: '{current_ai_q_text}', the user answered: '{user_answer_text_val}'. Does existing knowledge support this?"
                    relevant_indexed_chunks_for_validation = find_relevant_indexed_chunks(
                        rag_query,
                        KT_TARGET_RELEVANT_CONTENT_FILTER, # Changed
                        top_k=1
                    )
                if relevant_indexed_chunks_for_validation:
                    rag_chunk_for_val = relevant_indexed_chunks_for_validation[0]
                    raw_azure_search_score_from_rag = rag_chunk_for_val.get("@search.score", 0.0) # CAPTURE RAW AZURE SEARCH SCORE

                    rag_context_for_llm_val = rag_chunk_for_val.get('chunkText', '')
                    rag_source_file_for_llm_val = rag_chunk_for_val.get('sourceFileName', 'Indexed KB Document')

                    val_summary_rag, val_score_rag_llm_raw = validate_answer_with_llm(
                        current_ai_q_text, user_answer_text_val, rag_context_for_llm_val, f"RAG ({rag_source_file_for_llm_val})"
                    )
                    if val_score_rag_llm_raw is not None:
                        rag_llm_validation_score = val_score_rag_llm_raw
                        st.info(f"RAG Source: '{rag_source_file_for_llm_val}' (Raw Search Score: {raw_azure_search_score_from_rag:.2f}, LLM Validation Score: {rag_llm_validation_score:.2f})")

                        if rag_llm_validation_score > final_validation_score:
                            final_validation_score = rag_llm_validation_score
                            final_validation_source = f"RAG from KB: {rag_source_file_for_llm_val} (Search Score: {raw_azure_search_score_from_rag:.2f})"
                            validation_details_text = val_summary_rag
                    else:
                        st.warning(f"LLM validation of RAG chunk from '{rag_source_file_for_llm_val}' failed.")
                else:
                    st.info("No highly relevant RAG context found for this answer in the existing KB.")

                # 3. Conditionally run Tier 2 validation based on NEW LOGIC
                tier2_llm_validation_score = 0.0

                trigger_tier2_validation = False
                if raw_azure_search_score_from_rag is not None and raw_azure_search_score_from_rag >= 0.03: # Adjusted threshold, was 0.02 originally, check impact
                    st.info(f"RAG search score ({raw_azure_search_score_from_rag:.2f}) is high (>= 0.03). Tier 2 validation will be SKIPPED.")
                    trigger_tier2_validation = False
                elif final_validation_score < 0.6:
                    st.info(f"Current best validation score is {final_validation_score:.2f} (< 0.6), or RAG search score was low/None. Tier 2 validation may be triggered.")
                    trigger_tier2_validation = True
                else: # final_validation_score >= 0.6 AND (RAG raw score < 0.03 or RAG not used/found)
                    st.info(f"Current best validation score ({final_validation_score:.2f}) is adequate, and RAG search score (if any: {raw_azure_search_score_from_rag}) was not high enough to rely on solely or was low. Tier 2 not triggered by this path.")
                    trigger_tier2_validation = False

                if trigger_tier2_validation and st.session_state.all_blobs_for_kt_target_user: # Changed
                    st.info(f"Proceeding with Tier 2 deep dive on existing documents...")
                    with st.spinner("Performing Tier 2 validation deep dive (may take time)..."):
                        tier2_best_snippet_explanation, tier2_source_doc_filename, tier2_score_from_func = perform_tier2_validation_deep_dive(
                            current_ai_q_text, user_answer_text_val,
                            st.session_state.all_blobs_for_kt_target_user, # Changed
                            chat_client, blob_service_client, ADLS_CONTAINER_NAME
                        )
                    if tier2_score_from_func is not None and tier2_best_snippet_explanation:
                        tier2_llm_validation_score = tier2_score_from_func
                        st.info(f"Tier 2 validation found support in '{tier2_source_doc_filename}'. LLM Score: {tier2_llm_validation_score:.2f}")
                        if tier2_llm_validation_score > final_validation_score:
                            final_validation_score = tier2_llm_validation_score
                            final_validation_source = f"Tier 2 Deep Dive: {tier2_source_doc_filename}"
                            validation_details_text = tier2_best_snippet_explanation
                elif trigger_tier2_validation and not st.session_state.all_blobs_for_kt_target_user: # Changed
                    st.warning("Tier 2 validation was warranted based on scores, but no raw files are available for deep dive.")

                # Add Validation Observation to History
                validation_observation_content = {
                    "system_validation_event": "UserAnswerValidation",
                    "score": round(final_validation_score, 2),
                    "source": final_validation_source,
                    "details": validation_details_text[:500]
                }
                st.session_state.kt_interview_conversation_history.append({
                    "role": "system",
                    "content": json.dumps(validation_observation_content)
                })
                st.session_state.last_validation_display = {
                    "score": final_validation_score,
                    "source": final_validation_source,
                    "details": validation_details_text
                }

                # Document Ingestion (Conditional on final_validation_score)
                ingested_doc_msg = None # Initialize
                if current_q_processed_chunks and final_validation_score >= 0.7:
                    try:
                        with st.spinner(f"Ingesting content from '{current_q_uploaded_doc_name}' into Knowledge Base..."):
                            valid_chunks_to_ingest = [c for c in current_q_processed_chunks if c.get("vector") and c.get("id")]
                            if valid_chunks_to_ingest:
                                search_client.upload_documents(documents=valid_chunks_to_ingest)
                                ingested_doc_msg = f"Successfully ingested {len(valid_chunks_to_ingest)} chunks from '{current_q_uploaded_doc_name}'."
                                st.success(ingested_doc_msg)
                            else:
                                ingested_doc_msg = f"No valid chunks to ingest from '{current_q_uploaded_doc_name}' after processing."
                                st.warning(ingested_doc_msg)
                    except Exception as e_ingest_kt:
                        ingested_doc_msg = f"Error ingesting '{current_q_uploaded_doc_name}': {e_ingest_kt}"
                        st.error(ingested_doc_msg)
                elif current_q_processed_chunks:
                    ingested_doc_msg = f"Content from '{current_q_uploaded_doc_name}' was processed but not ingested due to validation score ({final_validation_score:.2f}) below threshold (0.7)."
                    st.warning(ingested_doc_msg)

                # Store Detailed Turn Result
                st.session_state.kt_interview_results.append({
                    "kt_session_id": st.session_state.kt_session_id,
                    "question_number": q_count,
                    "ai_question_details": dict(q_details),
                    "user_text_answer": user_answer_text_val,
                    "final_validation_source": final_validation_source,
                    "final_validation_details": validation_details_text,
                    "final_validation_score": final_validation_score,
                    "uploaded_document_name": current_q_uploaded_doc_name,
                    "ingestion_status_message": ingested_doc_msg,
                    "raw_azure_search_score_if_rag": raw_azure_search_score_from_rag,
                    "timestamp": format_date_for_azure_search(datetime.now(timezone.utc))
                })

            elif submit_kt_answer_button and not user_answer_text_val.strip():
                st.warning("Please provide an answer or type 'next' before submitting.")
                st.stop()

            # Generate Next AI Turn
            with st.spinner("CogniLink is formulating the next question..."):
                next_ai_question_details = generate_single_kt_interview_question(
                    employee_name=APP_KT_USER_NAME, # Changed
                    employee_email=APP_KT_USER_EMAIL, # Changed
                    conversation_history=st.session_state.kt_interview_conversation_history
                )

            st.session_state.current_kt_ai_question_details = next_ai_question_details
            if next_ai_question_details and "error" not in next_ai_question_details.get("question","").lower():
                st.session_state.kt_interview_conversation_history.append({
                    "role": "assistant",
                    "content": json.dumps(next_ai_question_details)
                })
                st.session_state.kt_interview_question_count += 1
            else:
                st.error("CogniLink had an issue generating the next question. The interview might conclude.")
                st.session_state.current_kt_ai_question_details = {"question": "[END_OF_INTERVIEW] Error occurred.", "response_structure":{}}

            st.rerun()

    # Interview Conclusion
    elif st.session_state.kt_interview_active and \
         (st.session_state.kt_interview_question_count > MAX_KT_INTERVIEW_QUESTIONS_TO_ASK or \
          (st.session_state.current_kt_ai_question_details and \
           ("[END_OF_INTERVIEW]" in st.session_state.current_kt_ai_question_details.get("question","") or \
            "error" in st.session_state.current_kt_ai_question_details.get("question","").lower()))):

        num_questions_actually_asked = st.session_state.kt_interview_question_count -1 if st.session_state.kt_interview_question_count > 0 else 0
        if "[END_OF_INTERVIEW]" not in st.session_state.current_kt_ai_question_details.get("question","") and \
           "error" not in st.session_state.current_kt_ai_question_details.get("question","").lower():
            st.info(f"Maximum number of interview turns ({MAX_KT_INTERVIEW_QUESTIONS_TO_ASK}) reached.")
            if st.session_state.current_kt_ai_question_details:
                 q_details_final = st.session_state.current_kt_ai_question_details
                 st.subheader(f"CogniLink's final intended question: \"{q_details_final['question']}\"")


        st.success(f"CogniLink Knowledge Transfer Interview Session Concluded. {num_questions_actually_asked} AI turns occurred (Configured max: {MAX_KT_INTERVIEW_QUESTIONS_TO_ASK}).")
        st.balloons()

        st.header("CogniLink KT Interview Summary")
        if st.session_state.kt_interview_results:
            for record in st.session_state.kt_interview_results:
                with st.expander(f"Turn {record['question_number']} (AI Q): {record['ai_question_details']['question'][:70]}...", expanded=False):
                    st.markdown(f"**CogniLink Question {record['question_number']}:** {record['ai_question_details']['question']}")
                    st.markdown(f"**User Answer:**")
                    st.text(record['user_text_answer'])
                    if record['uploaded_document_name']:
                        st.caption(f"Supporting Document: {record['uploaded_document_name']}")

                    score_color = "green" if record['final_validation_score'] >= 0.7 else "orange" if record['final_validation_score'] >= 0.4 else "red"
                    st.markdown(f"**Validation:** <span style='color:{score_color};'>Score `{record['final_validation_score']:.2f}`</span> from `{record['final_validation_source']}`", unsafe_allow_html=True)
                    st.caption(f"Details: *{record['final_validation_details']}*")
                    if record.get("raw_azure_search_score_if_rag") is not None:
                        st.caption(f"Raw Azure Search Score (if RAG used for validation): {record['raw_azure_search_score_if_rag']:.2f}")


                    if record['ingestion_status_message']:
                        if "Successfully ingested" in record['ingestion_status_message']:
                            st.success(f"KB Update: {record['ingestion_status_message']}")
                        elif "Error ingesting" in record['ingestion_status_message']:
                            st.error(f"KB Update: {record['ingestion_status_message']}")
                        else:
                            st.info(f"KB Info: {record['ingestion_status_message']}")
        else:
            st.write("No answers were recorded in this KT session.")

        if st.session_state.kt_interview_results or st.session_state.kt_interview_conversation_history:
            transcript_data = {
                "sessionId": st.session_state.kt_session_id,
                "employeeName": APP_KT_USER_NAME, # Changed
                "employeeEmail": APP_KT_USER_EMAIL, # Changed
                "maxTurnsConfigured": MAX_KT_INTERVIEW_QUESTIONS_TO_ASK,
                "actualTurns": num_questions_actually_asked,
                "interviewSummaryRecords": st.session_state.kt_interview_results,
                "fullConversationLog": st.session_state.kt_interview_conversation_history
            }
            try:
                transcript_json = json.dumps(transcript_data, indent=2)
                if st.download_button(
                    label="Download Full CogniLink Interview Transcript (JSON)",
                    data=transcript_json,
                    file_name=f"cognilink_kt_transcript_{st.session_state.kt_session_id}.json",
                    mime="application/json"
                ):
                    st.success("Transcript download initiated.")
            except TypeError as e_json:
                st.error(f"Error preparing transcript for download: {e_json}. Some data might not be serializable.")


    elif st.session_state.kt_interview_active:
        st.warning("CogniLink Interview is active but no question is currently loaded. Try starting a new session or check for errors.")

elif st.session_state.chat_mode:
    st.subheader(f"💬 Live Chat with {APP_KT_USER_NAME}'s Knowledge Base (Context-Aware)") # Changed

    for msg_idx, chat_message in enumerate(st.session_state.chat_history):
        with st.chat_message(chat_message["role"]):
            st.markdown(chat_message["content"])
            if chat_message["role"] == "assistant" and "sources" in chat_message and chat_message["sources"]:
                st.caption(f"Sources: {chat_message['sources']}")

    user_chat_input = st.chat_input(f"Ask a follow-up, or a new question about {APP_KT_USER_NAME}'s work...") # Changed

    if user_chat_input:
        st.session_state.chat_history.append({"role": "user", "content": user_chat_input})
        with st.chat_message("user"): st.markdown(user_chat_input)

        with st.chat_message("assistant"):
            final_chat_response = None; chat_response_tier = None; chat_sources_str = "N/A"

            COMMON_GREETINGS_AND_RESPONSES = {
                "hello": "Hello! How can I help you today?",
                "hi": "Hi there! What can I do for you?",
                "thanks": "You're welcome!",
                "thank you": "You're most welcome!"
            }
            chit_chat_reply = COMMON_GREETINGS_AND_RESPONSES.get(user_chat_input.lower().strip("!?."))

            if chit_chat_reply:
                final_chat_response = chit_chat_reply; chat_response_tier = "Conversational"
                st.markdown(final_chat_response)
            else:
                message_placeholder = st.empty(); message_placeholder.markdown("Thinking... 🤔")

                tier1_chat_answer, found_chat_ctx_t1, docs_t1 = answer_question_with_rag(
                    user_chat_input,
                    KT_TARGET_RELEVANT_CONTENT_FILTER, # Changed
                    top_k_contexts=3,
                    conversation_history=st.session_state.chat_history[:-1]
                )
                final_chat_response = tier1_chat_answer
                chat_response_tier = "Tier 1 (RAG, Context-Aware)"
                if docs_t1:
                    chat_sources_list = list(set(os.path.basename(d.get("sourceFileName", "N/A")) for d in docs_t1 if d.get("sourceFileName")))
                    chat_sources_str = ", ".join(chat_sources_list) if chat_sources_list else "N/A"

                if not found_chat_ctx_t1 or is_answer_unsatisfactory(tier1_chat_answer):
                    message_placeholder.markdown("Tier 1 answer seems limited. Trying Tier 2 (deep dive, context-aware)... 🕵️‍♂️")
                    if not st.session_state.all_blobs_for_kt_target_user: # Changed
                        final_chat_response = f"{tier1_chat_answer if found_chat_ctx_t1 else 'Information not found in indexed documents.'} Tier 2 deep dive is unavailable as no raw files are loaded."
                        chat_response_tier = "Tier 1 (No Tier 2 files)"
                    else:
                        t2_snippets = perform_tier2_deep_dive_for_one_question(
                            user_chat_input, st.session_state.all_blobs_for_kt_target_user, # Changed
                            chat_client, blob_service_client, ADLS_CONTAINER_NAME
                        )
                        if t2_snippets:
                            best_t2_ans, t2_src_ref = select_best_tier2_answer(
                                user_chat_input, t2_snippets, chat_client,
                                conversation_history=st.session_state.chat_history[:-1]
                            )
                            if best_t2_ans:
                                final_chat_response = best_t2_ans
                                chat_response_tier = "Tier 2 (Deep Dive from Raw Files, Context-Aware)"
                                chat_sources_str = t2_src_ref if t2_src_ref else "Raw Files (Tier 2)"

                message_placeholder.markdown(final_chat_response)
                if chat_sources_str and chat_sources_str != "N/A":
                    st.caption(f"Reference(s): {chat_sources_str}")

        st.session_state.chat_history.append({
            "role": "assistant",
            "content": final_chat_response,
            "tier": chat_response_tier,
            "sources": chat_sources_str
        })
        st.rerun()

else:
    if pipeline_choice == "Automated Q&A Pipeline":
        st.info("Select '🚀 Start Automated Knowledge Pipeline' from the sidebar.")
    elif pipeline_choice == "Structured KT Interview (CogniLink)":
        st.info(f"Select '🚀 Start New CogniLink KT Interview' with {APP_KT_USER_NAME} from the sidebar.") # Changed
    elif pipeline_choice == "Interactive Chat":
         st.info(f"Enter your question in the chat input field below to start a conversation about {APP_KT_USER_NAME}'s knowledge.") # Changed

# --- Debug Information (Optional Expander) ---
with st.sidebar.expander("🛠️ Debug Info", expanded=False):
    if st.session_state.debug_llm_messages_kt_qgen:
        st.write("Last LLM Messages for CogniLink QGen:")
        serializable_messages = []
        for msg in st.session_state.debug_llm_messages_kt_qgen:
            try:
                content = msg.get("content")
                if isinstance(content, str):
                    try:
                        parsed_content = json.loads(content)
                        serializable_messages.append({"role": msg.get("role"), "content": parsed_content})
                    except json.JSONDecodeError:
                        serializable_messages.append(msg)
                else:
                    serializable_messages.append(msg)
            except Exception:
                serializable_messages.append({"role": msg.get("role"), "content": str(msg.get("content"))})
        st.json(serializable_messages, expanded=False)

    if st.session_state.debug_kt_qgen_error_response:
        st.write("Last CogniLink QGen Error Response:")
        st.json(st.session_state.debug_kt_qgen_error_response, expanded=False)
    st.write("Current Session State Keys:", list(st.session_state.keys()))
    st.write(f"Max KT Questions (from session_state): {st.session_state.get('max_kt_questions_user_set', 'Not Set')}")
    if st.session_state.kt_interview_conversation_history:
        st.write("KT Interview Conversation History (last 5 turns):")
        serializable_conv_hist = []
        for turn in st.session_state.kt_interview_conversation_history[-5:]:
            try:
                content = turn.get("content")
                if isinstance(content, str):
                    try:
                        parsed_content = json.loads(content)
                        serializable_conv_hist.append({"role": turn.get("role"), "content": parsed_content})
                    except json.JSONDecodeError:
                         serializable_conv_hist.append(turn)
                else:
                    serializable_conv_hist.append(turn)
            except Exception:
                 serializable_conv_hist.append({"role": turn.get("role"), "content": str(turn.get("content"))})

        st.json(serializable_conv_hist, expanded=False)

